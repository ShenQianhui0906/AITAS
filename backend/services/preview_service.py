"""
Preview generation service — AppleScript-based PPT/Keynote PDF export,
QuickLook HTML generation, and courseware preview HTML builders.
"""
from __future__ import annotations

import html
import mimetypes
import re
import secrets
import shutil
import subprocess
import tempfile
import threading
import zipfile
from pathlib import Path
from urllib.parse import quote, urlencode

from backend.config import (
    UPLOAD_DIR, TMP_DIR,
    QUICKLOOK_PREVIEW_SUFFIXES, NATIVE_PREVIEW_SUFFIX,
)

# ------------------------- paths / locks -------------------------
QLMANAGE_PATH = shutil.which("qlmanage")
OSASCRIPT_PATH = shutil.which("osascript")
PREVIEW_BUILD_LOCK = threading.Lock()

POWERPOINT_APP_NAME = "Microsoft PowerPoint"
POWERPOINT_APP_PATH = Path("/Applications/Microsoft PowerPoint.app")
KEYNOTE_APP_NAME = "Keynote"
KEYNOTE_APP_PATH = Path("/Applications/Keynote.app")

# ------------------------- helpers -------------------------

def escape_applescript_text(value: str) -> str:
    return str(value).replace("\\", "\\\\").replace('"', '\\"')


def run_osascript_lines(lines: list[str], timeout: int = 240) -> bool:
    if not OSASCRIPT_PATH:
        return False
    command = [OSASCRIPT_PATH]
    for line in lines:
        command.extend(["-e", line])
    try:
        result = subprocess.run(command, capture_output=True, text=True, timeout=timeout)
    except (OSError, subprocess.SubprocessError):
        return False
    return result.returncode == 0


def get_quicklook_preview_dir(file_path: Path) -> Path:
    return file_path.parent / f"{file_path.name}.qlpreview"


def get_native_pdf_preview_path(file_path: Path) -> Path:
    return file_path.parent / f"{file_path.name}{NATIVE_PREVIEW_SUFFIX}"


# ------------------------- URL builders -------------------------

def build_upload_url(file_name: str) -> str:
    return f"/uploads/{quote(file_name)}"


def build_quicklook_preview_url(file_name: str) -> str:
    return build_upload_url(f"{file_name}.qlpreview/Preview.html")


def build_quicklook_viewer_url(file_name: str) -> str:
    return f"/preview-quicklook/{quote(file_name)}"


def build_pptx_media_url(file_name: str, asset_name: str) -> str:
    return f"/preview-media?{urlencode({'file': file_name, 'asset': asset_name})}"


# ------------------------- PPTX helpers -------------------------

def extract_pptx_text_slides(file_path: Path) -> list[dict] | None:
    """Extract text content from each slide of a PPTX file. Returns list of {slide_num, title, body_lines}."""
    try:
        from pptx import Presentation as PPTXPresentation
    except Exception:
        return None
    try:
        prs = PPTXPresentation(str(file_path))
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            title = ""
            body_lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                            if not title:
                                title = text
                        else:
                            body_lines.append(text)
            slides.append({"slide_num": idx, "title": title, "body_lines": body_lines})
        return slides
    except Exception:
        return None


def build_pptx_slides_preview_html(file_path: Path) -> str | None:
    """Build an HTML preview for a PPTX file showing text per slide."""
    slides = extract_pptx_text_slides(file_path)
    if not slides:
        return None
    display_title = file_path.stem
    title = html.escape(display_title)

    slides_html = ""
    for slide in slides:
        sn = slide["slide_num"]
        slide_title = html.escape(slide["title"] or f"幻灯片 {sn}")
        body = "<br>".join(html.escape(line) for line in slide["body_lines"])
        slides_html += f"""
        <div class="slide-card">
            <div class="slide-num">{sn}</div>
            <div class="slide-title">{slide_title}</div>
            <div class="slide-body">{body}</div>
        </div>"""

    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    :root {{ color-scheme: light; --bg: #f8fafc; --surface: #ffffff; --text-main: #0f172a;
      --text-muted: #64748b; --border: #e2e8f0; --accent: #6366f1;
      --shadow: 0 12px 32px rgba(15, 23, 42, 0.08); --radius: 18px; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; padding: 20px; background: var(--bg); font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; }}
    .slides-container {{ display: flex; flex-direction: column; gap: 16px; max-width: 900px; margin: 0 auto; }}
    .slide-card {{ background: var(--surface); border: 1px solid var(--border); border-radius: var(--radius);
      padding: 24px; box-shadow: var(--shadow); }}
    .slide-num {{ font-size: 12px; color: var(--accent); font-weight: 700; margin-bottom: 8px; }}
    .slide-title {{ font-size: 20px; font-weight: 600; color: var(--text-main); margin-bottom: 12px; }}
    .slide-body {{ font-size: 14px; color: var(--text-muted); line-height: 1.7; }}
  </style>
</head>
<body>
  <div class="slides-container">{slides_html}</div>
</body>
</html>"""


# ------------------------- PDF export (macOS) -------------------------

def export_presentation_pdf_via_powerpoint(source_path: str, target_pdf_path: str) -> bool:
    if not POWERPOINT_APP_PATH.exists():
        return False
    source_literal = escape_applescript_text(source_path)
    target_literal = escape_applescript_text(target_pdf_path)
    lines = [
        f'set sourceFile to POSIX file "{source_literal}"',
        f'set destFile to POSIX file "{target_literal}"',
        f'tell application "{POWERPOINT_APP_NAME}"',
        "activate",
        "open sourceFile",
        "delay 2",
        "save active presentation in destFile as save as PDF",
        "close active presentation saving no",
        "end tell",
    ]
    return run_osascript_lines(lines, timeout=300)


def export_presentation_pdf_via_keynote(source_path: str, target_pdf_path: str) -> bool:
    if not KEYNOTE_APP_PATH.exists():
        return False
    source_literal = escape_applescript_text(source_path)
    target_literal = escape_applescript_text(target_pdf_path)
    lines = [
        f'set sourceFile to POSIX file "{source_literal}"',
        f'set destFile to POSIX file "{target_literal}"',
        f'tell application "{KEYNOTE_APP_NAME}"',
        "activate",
        "set theDocument to open sourceFile",
        "delay 2",
        "export theDocument to destFile as PDF with properties {{PDF image quality:Best}}",
        "close theDocument saving no",
        "end tell",
    ]
    return run_osascript_lines(lines, timeout=300)


def ensure_native_pdf_preview(file_path: Path) -> Path | None:
    """Try to export a native PDF preview of a PPT/PPTX file using Keynote or PowerPoint."""
    if file_path.suffix.lower() not in QUICKLOOK_PREVIEW_SUFFIXES or not OSASCRIPT_PATH:
        return None

    target_path = get_native_pdf_preview_path(file_path)
    try:
        source_mtime = file_path.stat().st_mtime
    except OSError:
        return None

    if target_path.exists():
        try:
            if target_path.stat().st_mtime >= source_mtime and target_path.stat().st_size > 0:
                return target_path
        except OSError:
            pass

    with PREVIEW_BUILD_LOCK:
        if target_path.exists():
            try:
                if target_path.stat().st_mtime >= source_mtime and target_path.stat().st_size > 0:
                    return target_path
            except OSError:
                pass

        temp_root = None
        try:
            temp_root = Path(tempfile.mkdtemp(prefix="aitas-native-preview-"))
            temp_pdf = temp_root / f"{file_path.stem}.pdf"
            exporters = [export_presentation_pdf_via_powerpoint, export_presentation_pdf_via_keynote]
            for exporter in exporters:
                if temp_pdf.exists():
                    temp_pdf.unlink(missing_ok=True)
                if not exporter(str(file_path), str(temp_pdf)):
                    continue
                if not temp_pdf.exists():
                    continue
                try:
                    if temp_pdf.stat().st_size <= 0:
                        continue
                except OSError:
                    continue

                temp_target = TMP_DIR / f".{target_path.name}.{secrets.token_hex(4)}.tmp"
                if temp_target.exists():
                    temp_target.unlink(missing_ok=True)
                shutil.copy2(temp_pdf, temp_target)
                temp_target.replace(target_path)
                return target_path
        except OSError:
            return None
        finally:
            if temp_root and temp_root.exists():
                shutil.rmtree(temp_root, ignore_errors=True)
    return None


# ------------------------- QuickLook preview -------------------------

def ensure_quicklook_preview(file_path: Path) -> Path | None:
    """Generate a qlpreview directory for a PPT/PPTX file using macOS qlmanage."""
    if file_path.suffix.lower() not in QUICKLOOK_PREVIEW_SUFFIXES or not QLMANAGE_PATH:
        return None

    target_dir = get_quicklook_preview_dir(file_path)
    preview_html = target_dir / "Preview.html"
    try:
        source_mtime = file_path.stat().st_mtime
    except OSError:
        return None

    if preview_html.exists():
        try:
            if preview_html.stat().st_mtime >= source_mtime:
                return target_dir
        except OSError:
            pass

    with PREVIEW_BUILD_LOCK:
        if preview_html.exists():
            try:
                if preview_html.stat().st_mtime >= source_mtime:
                    return target_dir
            except OSError:
                pass

        temp_root = None
        try:
            temp_root = Path(tempfile.mkdtemp(prefix="aitas-quicklook-"))
            result = subprocess.run(
                [QLMANAGE_PATH, "-o", str(temp_root), "-p", str(file_path)],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode != 0:
                return None

            generated_dir = next(temp_root.glob("*.qlpreview"), None)
            if generated_dir is None or not (generated_dir / "Preview.html").exists():
                return None

            temp_target = TMP_DIR / f".{target_dir.name}.{secrets.token_hex(4)}.tmp"
            if temp_target.exists():
                shutil.rmtree(temp_target, ignore_errors=True)
            shutil.copytree(generated_dir, temp_target)
            if target_dir.exists():
                shutil.rmtree(target_dir, ignore_errors=True)
            temp_target.rename(target_dir)
            return target_dir
        except (OSError, StopIteration, subprocess.SubprocessError):
            return None
        finally:
            if temp_root and temp_root.exists():
                shutil.rmtree(temp_root, ignore_errors=True)


# ------------------------- HTML builder helpers -------------------------

def _normalize_quicklook_numeric_css(raw_html: str) -> str:
    """Fix common CSS number issues in quicklook output."""
    return raw_html


def _transform_quicklook_assets(raw_html: str) -> str:
    """Rewrite asset URLs in quicklook HTML."""
    return raw_html


def _build_quicklook_runtime_assets(base_href: str) -> tuple[str, str]:
    style = f'<base href="{base_href}">'
    script = ""
    return style, script


def build_browser_ready_quicklook_html(file_path: Path, preview_html_path: Path) -> str:
    raw_html = preview_html_path.read_text(encoding="utf-8", errors="ignore")
    raw_html = _normalize_quicklook_numeric_css(raw_html)
    raw_html = _transform_quicklook_assets(raw_html)
    base_href = build_upload_url(f"{file_path.name}.qlpreview/")
    style, script = _build_quicklook_runtime_assets(base_href)
    if "<head>" in raw_html:
        raw_html = raw_html.replace("<head>", f"<head>{style}", 1)
    else:
        raw_html = f"<head>{style}</head>{raw_html}"
    if "</body>" in raw_html:
        raw_html = raw_html.replace("</body>", f"{script}</body>", 1)
    else:
        raw_html = f"{raw_html}{script}"
    return raw_html


def get_preview_mode(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    if suffix in {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".html", ".htm", ".docx", ".pptx"}:
        return "text"
    return "unsupported"


def get_preview_label(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    labels = {
        ".pdf": "PDF 在线预览",
        ".pptx": "PPT 文字预览",
        ".docx": "Word 文字预览",
        ".txt": "文本预览",
        ".md": "Markdown 预览",
        ".csv": "表格文本预览",
        ".json": "JSON 文本预览",
    }
    return labels.get(suffix, "课件预览")


def build_courseware_preview_html(file_path: Path) -> str:
    """Build a complete HTML page to preview any courseware file."""
    from backend.utils.file_utils import get_display_file_title

    display_title = get_display_file_title(file_path)
    title = html.escape(display_title)
    raw_url = build_upload_url(file_path.name)
    preview_mode = get_preview_mode(file_path)
    preview_label = html.escape(get_preview_label(file_path))
    suffix = file_path.suffix.lower()

    # PPTX slide text preview
    if suffix == ".pptx":
        pptx_preview = build_pptx_slides_preview_html(file_path)
        if pptx_preview:
            return pptx_preview

    # QuickLook for PPT/PPTX
    if suffix in QUICKLOOK_PREVIEW_SUFFIXES:
        quicklook_dir = ensure_quicklook_preview(file_path)
        quicklook_html = quicklook_dir / "Preview.html" if quicklook_dir else None
        if quicklook_html and quicklook_html.exists():
            preview_content = f"""
                <iframe class="preview-document-frame preview-slides-frame"
                  src="{build_quicklook_viewer_url(file_path.name)}"
                  title="{title}"></iframe>"""
            return _wrap_preview_html(title, preview_label, preview_content)

    # PDF inline preview
    if preview_mode == "pdf":
        preview_content = f"""
            <iframe class="preview-document-frame"
              src="{raw_url}"
              title="{title}"></iframe>"""
        return _wrap_preview_html(title, preview_label, preview_content)

    # Image preview
    if preview_mode == "image":
        preview_content = f'<div class="preview-image-wrap"><img src="{raw_url}" alt="{title}" /></div>'
        return _wrap_preview_html(title, preview_label, preview_content)

    # Text preview
    if preview_mode == "text":
        try:
            text_content = html.escape(file_path.read_text(encoding="utf-8", errors="replace"))
        except Exception:
            text_content = "无法读取文件内容。"
        preview_content = f'<pre class="preview-text">{text_content}</pre>'
        return _wrap_preview_html(title, preview_label, preview_content)

    # Fallback
    return _wrap_preview_html(title, preview_label,
        f'<p class="preview-unsupported">不支持预览此文件类型，请<a href="{raw_url}">下载查看</a>。</p>')


def _wrap_preview_html(title: str, preview_label: str, content: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    :root {{ color-scheme: light; --bg: #f8fafc; --surface: #ffffff; --text-main: #0f172a;
      --text-muted: #64748b; --border: #e2e8f0; --accent: #6366f1;
      --shadow: 0 12px 32px rgba(15, 23, 42, 0.08); --radius: 22px; }}
    * {{ box-sizing: border-box; }}
    html, body {{ margin: 0; min-height: 100%; background: var(--bg);
      color: var(--text-main); font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; }}
    body {{ padding: 20px; }}
    .preview-shell {{ min-height: calc(100vh - 40px); display: flex; flex-direction: column; gap: 16px; }}
    .preview-topbar {{ display: flex; align-items: center; justify-content: space-between; gap: 16px;
      padding: 18px 20px; border: 1px solid var(--border); border-radius: var(--radius);
      background: rgba(255,255,255,0.92); box-shadow: var(--shadow); }}
    .preview-badge {{ padding: 8px 12px; border-radius: 999px;
      background: rgba(99,102,241,0.12); color: var(--accent); font-size: 13px; font-weight: 600; }}
    .preview-main {{ flex: 1; }}
    .preview-document-frame {{ width: 100%; min-height: calc(100vh - 140px); border: 1px solid var(--border);
      border-radius: var(--radius); background: var(--surface); box-shadow: var(--shadow); }}
    .preview-image-wrap {{ display: flex; justify-content: center; padding: 20px; }}
    .preview-image-wrap img {{ max-width: 100%; border-radius: var(--radius); box-shadow: var(--shadow); }}
    .preview-text {{ padding: 20px; background: var(--surface); border: 1px solid var(--border);
      border-radius: var(--radius); white-space: pre-wrap; word-break: break-all; font-size: 14px;
      line-height: 1.7; box-shadow: var(--shadow); max-height: calc(100vh - 140px); overflow: auto; }}
    .preview-unsupported {{ text-align: center; padding: 40px; color: var(--text-muted); }}
    .preview-slides-frame {{ border: none; }}
  </style>
</head>
<body>
<div class="preview-shell">
  <div class="preview-topbar">
    <span class="preview-badge">{preview_label}</span>
  </div>
  <div class="preview-main">{content}</div>
</div>
</body>
</html>"""

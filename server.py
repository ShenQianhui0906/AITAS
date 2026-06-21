import argparse
import html
import hashlib
import json
import mimetypes
import os
import re
import secrets
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import threading
import time
import urllib.error
import urllib.request
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from email.parser import BytesParser
from email.policy import default
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import parse_qs, quote, unquote, urlencode, urlparse


ROOT_DIR = Path(__file__).resolve().parent
STATIC_DIR = ROOT_DIR / "static"
UPLOAD_DIR = ROOT_DIR / "uploads"
DATA_DIR = ROOT_DIR / "data"
DB_PATH = DATA_DIR / "ai_tutor.db"
VENDOR_DIR = ROOT_DIR / "vendor"
if VENDOR_DIR.exists():
    sys.path.insert(0, str(VENDOR_DIR))

_ENV_FILE = ROOT_DIR / "ENV"
if _ENV_FILE.exists():
    with _ENV_FILE.open(encoding="utf-8") as _f:
        for _line in _f:
            _line = _line.strip()
            if not _line or _line.startswith("#") or "=" not in _line:
                continue
            _key, _, _val = _line.partition("=")
            os.environ.setdefault(_key.strip(), _val.strip())

try:
    import rag as _rag_module

    RAG_AVAILABLE = True
except ImportError:
    _rag_module = None  # type: ignore[assignment]
    RAG_AVAILABLE = False

ADMIN_USERNAME = "admin"
ADMIN_PASSWORD = "2026"
ADMIN_DISPLAY_NAME = "系统管理员"
BIGMODEL_API_URL = os.environ.get("BIGMODEL_API_URL", "https://open.bigmodel.cn/api/paas/v4/chat/completions").strip() or "https://open.bigmodel.cn/api/paas/v4/chat/completions"
BIGMODEL_MODEL = os.environ.get("BIGMODEL_MODEL", "glm-4.7-flash").strip() or "glm-4.7-flash"
MAX_AI_CONTEXT_CHARS = 12000
MAX_AI_HISTORY_MESSAGES = 8
SYNC_CONDITION = threading.Condition()
SYNC_CURSOR = 0
USER_SYNC_CURSORS = {}
QLMANAGE_PATH = shutil.which("qlmanage")
OSASCRIPT_PATH = shutil.which("osascript")
QUICKLOOK_PREVIEW_SUFFIXES = {".ppt", ".pptx"}
PREVIEW_BUILD_LOCK = threading.Lock()
NATIVE_PREVIEW_SUFFIX = ".native-preview.pdf"
POWERPOINT_APP_NAME = "Microsoft PowerPoint"
POWERPOINT_APP_PATH = Path("/Applications/Microsoft PowerPoint.app")
KEYNOTE_APP_NAME = "Keynote"
KEYNOTE_APP_PATH = Path("/Applications/Keynote.app")

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None


def now_iso():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def build_upload_url(file_name):
    return f"/uploads/{quote(file_name)}"


def get_preview_display_name(file_name, display_title=None):
    suffix = Path(file_name).suffix
    title = (display_title or "").strip() or get_display_file_title(Path(file_name))
    title = re.sub(r'[\\/:*?"<>|]+', " ", title).strip().rstrip(".")
    title = re.sub(r"\s+", " ", title) or get_display_file_title(Path(file_name)) or "课件"
    if suffix and not title.lower().endswith(suffix.lower()):
        title = f"{title}{suffix}"
    return title


def build_viewer_url(file_name, display_title=None):
    preview_name = get_preview_display_name(file_name, display_title)
    return f"/preview/{quote(file_name)}/{quote(preview_name)}"


def build_quicklook_preview_url(file_name):
    return build_upload_url(f"{file_name}.qlpreview/Preview.html")


def build_quicklook_viewer_url(file_name):
    return f"/preview-quicklook/{quote(file_name)}"


def build_pptx_media_url(file_name, asset_name):
    return f"/preview-media?{urlencode({'file': file_name, 'asset': asset_name})}"


def get_quicklook_preview_dir(file_path):
    return UPLOAD_DIR / f"{file_path.name}.qlpreview"


def get_native_pdf_preview_path(file_path):
    return UPLOAD_DIR / f"{file_path.name}{NATIVE_PREVIEW_SUFFIX}"


def escape_applescript_text(value):
    return str(value).replace("\\", "\\\\").replace('"', '\\"')


def run_osascript_lines(lines, timeout=240):
    if not OSASCRIPT_PATH:
        return False
    command = [OSASCRIPT_PATH]
    for line in lines:
        command.extend(["-e", line])
    try:
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
    except (OSError, subprocess.SubprocessError):
        return False
    return result.returncode == 0


def export_presentation_pdf_via_powerpoint(source_path, target_pdf_path):
    if not POWERPOINT_APP_PATH.exists():
        return False
    source_literal = escape_applescript_text(str(source_path))
    target_literal = escape_applescript_text(str(target_pdf_path))
    lines = [
        f'set sourceFile to POSIX file "{source_literal}"',
        f'set destFile to POSIX file "{target_literal}"',
        f'tell application "{POWERPOINT_APP_NAME}"',
        "activate",
        "open sourceFile",
        "delay 2",
        "save active presentation in destFile as save as PDF",
        "close active presentation saving no",
        "end tell",
    ]
    return run_osascript_lines(lines, timeout=300)


def export_presentation_pdf_via_keynote(source_path, target_pdf_path):
    if not KEYNOTE_APP_PATH.exists():
        return False
    source_literal = escape_applescript_text(str(source_path))
    target_literal = escape_applescript_text(str(target_pdf_path))
    lines = [
        f'set sourceFile to POSIX file "{source_literal}"',
        f'set destFile to POSIX file "{target_literal}"',
        f'tell application "{KEYNOTE_APP_NAME}"',
        "activate",
        "set theDocument to open sourceFile",
        "delay 2",
        "export theDocument to destFile as PDF with properties {PDF image quality:Best}",
        "close theDocument saving no",
        "end tell",
    ]
    return run_osascript_lines(lines, timeout=300)


def ensure_native_pdf_preview(file_path):
    if file_path.suffix.lower() not in QUICKLOOK_PREVIEW_SUFFIXES or not OSASCRIPT_PATH:
        return None

    target_path = get_native_pdf_preview_path(file_path)
    try:
        source_mtime = file_path.stat().st_mtime
    except OSError:
        return None

    if target_path.exists():
        try:
            if target_path.stat().st_mtime >= source_mtime and target_path.stat().st_size > 0:
                return target_path
        except OSError:
            pass

    with PREVIEW_BUILD_LOCK:
        if target_path.exists():
            try:
                if target_path.stat().st_mtime >= source_mtime and target_path.stat().st_size > 0:
                    return target_path
            except OSError:
                pass

        temp_root = None
        try:
            temp_root = Path(tempfile.mkdtemp(prefix="aitas-native-preview-"))
            temp_pdf = temp_root / f"{file_path.stem}.pdf"
            exporters = [
                export_presentation_pdf_via_powerpoint,
                export_presentation_pdf_via_keynote,
            ]
            for exporter in exporters:
                if temp_pdf.exists():
                    temp_pdf.unlink(missing_ok=True)
                if not exporter(file_path, temp_pdf):
                    continue
                if not temp_pdf.exists():
                    continue
                try:
                    if temp_pdf.stat().st_size <= 0:
                        continue
                except OSError:
                    continue

                temp_target = UPLOAD_DIR / f".{target_path.name}.{secrets.token_hex(4)}.tmp"
                if temp_target.exists():
                    temp_target.unlink(missing_ok=True)
                shutil.copy2(temp_pdf, temp_target)
                temp_target.replace(target_path)
                return target_path
        except OSError:
            return None
        finally:
            if temp_root and temp_root.exists():
                shutil.rmtree(temp_root, ignore_errors=True)
    return None


def ensure_quicklook_preview(file_path):
    if file_path.suffix.lower() not in QUICKLOOK_PREVIEW_SUFFIXES or not QLMANAGE_PATH:
        return None

    target_dir = get_quicklook_preview_dir(file_path)
    preview_html = target_dir / "Preview.html"
    try:
        source_mtime = file_path.stat().st_mtime
    except OSError:
        return None

    if preview_html.exists():
        try:
            if preview_html.stat().st_mtime >= source_mtime:
                return target_dir
        except OSError:
            pass

    with PREVIEW_BUILD_LOCK:
        if preview_html.exists():
            try:
                if preview_html.stat().st_mtime >= source_mtime:
                    return target_dir
            except OSError:
                pass

        temp_root = None
        try:
            temp_root = Path(tempfile.mkdtemp(prefix="aitas-quicklook-"))
            result = subprocess.run(
                [QLMANAGE_PATH, "-o", str(temp_root), "-p", str(file_path)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                return None

            generated_dir = next(temp_root.glob("*.qlpreview"), None)
            if generated_dir is None or not (generated_dir / "Preview.html").exists():
                return None

            temp_target = UPLOAD_DIR / f".{target_dir.name}.{secrets.token_hex(4)}.tmp"
            if temp_target.exists():
                shutil.rmtree(temp_target, ignore_errors=True)
            shutil.copytree(generated_dir, temp_target)
            if target_dir.exists():
                shutil.rmtree(target_dir, ignore_errors=True)
            temp_target.rename(target_dir)
            return target_dir
        except (OSError, StopIteration, subprocess.SubprocessError):
            return None
        finally:
            if temp_root and temp_root.exists():
                shutil.rmtree(temp_root, ignore_errors=True)


def delete_courseware_assets(stored_file_name):
    if not stored_file_name or stored_file_name == "demo_courseware.txt":
        return

    file_path = UPLOAD_DIR / stored_file_name
    if file_path.exists():
        if file_path.is_dir():
            shutil.rmtree(file_path, ignore_errors=True)
        else:
            file_path.unlink()

    preview_dir = UPLOAD_DIR / f"{stored_file_name}.qlpreview"
    if preview_dir.exists():
        shutil.rmtree(preview_dir, ignore_errors=True)

    native_preview = get_native_pdf_preview_path(UPLOAD_DIR / stored_file_name)
    if native_preview.exists():
        native_preview.unlink(missing_ok=True)


def get_display_file_title(file_path):
    name = file_path.stem
    name = re.sub(r"^[0-9a-f]{12,}_", "", name, flags=re.IGNORECASE)
    name = re.sub(r"^\d+[._-]+", "", name)
    return name or file_path.stem


def emu_to_px(value):
    return round(float(value) / 9525.0, 2)


def pt_to_px(value):
    return round(float(value) * 96.0 / 72.0, 2)


def css_color_from_rgb(rgb_value):
    if not rgb_value:
        return None
    rgb_text = str(rgb_value).strip()
    return f"#{rgb_text}" if rgb_text else None


def get_shape_fill_css(shape):
    fill = getattr(shape, "fill", None)
    if not fill:
        return ""
    try:
        color = css_color_from_rgb(fill.fore_color.rgb)
    except Exception:
        color = None
    if color:
        return f"background:{color};"
    return ""


def get_shape_line_css(shape):
    line = getattr(shape, "line", None)
    if not line:
        return ""
    try:
        color = css_color_from_rgb(line.color.rgb)
    except Exception:
        color = None
    width = None
    try:
        width = emu_to_px(line.width) if line.width else None
    except Exception:
        width = None
    if color:
        return f"border:{max(width or 1.0, 1.0)}px solid {color};"
    return ""


def get_shape_transform(shape, parent_transform):
    left = parent_transform["origin_x"] + (shape.left - parent_transform["child_origin_x"]) * parent_transform["scale_x"]
    top = parent_transform["origin_y"] + (shape.top - parent_transform["child_origin_y"]) * parent_transform["scale_y"]
    width = shape.width * parent_transform["scale_x"]
    height = shape.height * parent_transform["scale_y"]
    return left, top, width, height


def build_group_transform(shape, parent_transform):
    left, top, width, height = get_shape_transform(shape, parent_transform)
    xfrm = getattr(shape.element, "xfrm", None)
    child_origin_x = 0
    child_origin_y = 0
    scale_x = parent_transform["scale_x"]
    scale_y = parent_transform["scale_y"]
    if xfrm is not None:
        try:
            child_origin_x = xfrm.chOff.x
            child_origin_y = xfrm.chOff.y
            if xfrm.chExt.cx:
                scale_x = parent_transform["scale_x"] * (shape.width / xfrm.chExt.cx)
            if xfrm.chExt.cy:
                scale_y = parent_transform["scale_y"] * (shape.height / xfrm.chExt.cy)
        except Exception:
            child_origin_x = 0
            child_origin_y = 0
            scale_x = parent_transform["scale_x"]
            scale_y = parent_transform["scale_y"]
    return {
        "origin_x": left,
        "origin_y": top,
        "child_origin_x": child_origin_x,
        "child_origin_y": child_origin_y,
        "scale_x": scale_x,
        "scale_y": scale_y,
    }


def shape_text_payload(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    text_frame = shape.text_frame
    paragraphs = []
    sample_font_px = None
    sample_color = None
    sample_bold = None
    sample_align = None
    for paragraph in text_frame.paragraphs:
        text = paragraph.text or ""
        text = text.replace("\x0b", "\n").strip()
        if not text:
            continue
        if sample_align is None and paragraph.alignment is not None:
            sample_align = str(paragraph.alignment).split(" ")[0].lower()
        for run in paragraph.runs:
            if sample_font_px is None and run.font.size:
                sample_font_px = pt_to_px(run.font.size.pt)
            if sample_color is None:
                try:
                    sample_color = css_color_from_rgb(run.font.color.rgb)
                except Exception:
                    sample_color = None
            if sample_bold is None and run.font.bold is not None:
                sample_bold = bool(run.font.bold)
        paragraphs.append({"text": text, "level": getattr(paragraph, "level", 0) or 0})
    if not paragraphs:
        return None
    return {
        "paragraphs": paragraphs,
        "font_px": sample_font_px or 22,
        "color": sample_color or "#0f172a",
        "bold": bool(sample_bold),
        "align": sample_align or "left",
        "vertical_anchor": str(getattr(text_frame, "vertical_anchor", "")).lower(),
    }


def table_payload(shape):
    if not getattr(shape, "has_table", False):
        return None
    rows = []
    for row in shape.table.rows:
        cells = []
        for cell in row.cells:
            text = normalize_text_content(cell.text or "")
            cells.append(text)
        if any(cells):
            rows.append(cells)
    return rows or None


def picture_payload(slide, shape):
    if MSO_SHAPE_TYPE is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        rel_id = shape._pic.blipFill.blip.rEmbed
        part = slide.part.related_part(rel_id)
        asset_name = Path(str(part.partname)).name
    except Exception:
        return None
    ext = Path(asset_name).suffix.lower()
    if ext not in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}:
        return None
    return {
        "asset_name": asset_name,
        "alt": html.escape(shape.name or "slide image"),
    }


def collect_pptx_slide_items(slide, shapes, transform, items):
    for shape in shapes:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            collect_pptx_slide_items(slide, shape.shapes, build_group_transform(shape, transform), items)
            continue

        left, top, width, height = get_shape_transform(shape, transform)
        if width <= 1 or height <= 1:
            continue

        picture = picture_payload(slide, shape)
        if picture:
            items.append(
                {
                    "kind": "image",
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "asset_name": picture["asset_name"],
                    "alt": picture["alt"],
                }
            )
            continue

        table_rows = table_payload(shape)
        if table_rows:
            items.append(
                {
                    "kind": "table",
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "rows": table_rows,
                }
            )
            continue

        text = shape_text_payload(shape)
        if text:
            items.append(
                {
                    "kind": "text",
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "payload": text,
                    "fill_css": get_shape_fill_css(shape),
                    "line_css": get_shape_line_css(shape),
                    "shape_type": str(shape.shape_type),
                }
            )
            continue

        fill_css = get_shape_fill_css(shape)
        if fill_css and width * height > 6400:
            items.append(
                {
                    "kind": "box",
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "fill_css": fill_css,
                    "line_css": get_shape_line_css(shape),
                }
            )


def detect_slide_title(slide, items):
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None:
        title_text = normalize_text_content(getattr(title_shape, "text", ""))
        if title_text:
            return title_text

    text_items = [item for item in items if item["kind"] == "text"]
    if not text_items:
        return ""
    top_item = min(text_items, key=lambda item: (item["top"], -item["payload"]["font_px"]))
    first_text = top_item["payload"]["paragraphs"][0]["text"]
    return normalize_text_content(first_text)


def render_pptx_text_html(payload):
    parts = []
    for paragraph in payload["paragraphs"]:
        text = html.escape(paragraph["text"]).replace("\n", "<br />")
        indent = paragraph["level"] * 1.25
        parts.append(f'<p style="margin:0 0 0.35em;padding-left:{indent}em;">{text}</p>')
    return "".join(parts)


def render_pptx_slide_item(file_path, item):
    left_px = emu_to_px(item["left"])
    top_px = emu_to_px(item["top"])
    width_px = emu_to_px(item["width"])
    height_px = emu_to_px(item["height"])
    style = (
        f"left:{left_px:.2f}px;"
        f"top:{top_px:.2f}px;"
        f"width:{width_px:.2f}px;"
        f"height:{height_px:.2f}px;"
    )
    if item["kind"] == "image":
        src = build_pptx_media_url(file_path.name, item["asset_name"])
        return (
            f'<img class="ppt-render-image" src="{src}" alt="{item["alt"]}" '
            f'style="position:absolute;{style}" loading="lazy" />'
        )
    if item["kind"] == "table":
        rows_html = []
        for row in item["rows"]:
            cells = "".join(f"<td>{html.escape(cell)}</td>" for cell in row)
            rows_html.append(f"<tr>{cells}</tr>")
        return (
            f'<div class="ppt-render-table-wrap" style="position:absolute;{style}">'
            f'<table class="ppt-render-table">{"".join(rows_html)}</table>'
            f"</div>"
        )
    if item["kind"] == "box":
        return (
            f'<div class="ppt-render-box" style="position:absolute;{style}{item["fill_css"]}{item["line_css"]}"></div>'
        )
    payload = item["payload"]
    text_style = (
        f"position:absolute;{style}"
        f"font-size:{payload['font_px']:.2f}px;"
        f"color:{payload['color']};"
        f"text-align:{payload['align']};"
        f"font-weight:{700 if payload['bold'] else 500};"
        f"{item['fill_css']}{item['line_css']}"
    )
    padding = "padding:10px 14px;" if item["fill_css"] else ""
    return (
        f'<div class="ppt-render-text" style="{text_style}{padding}">'
        f'{render_pptx_text_html(payload)}'
        f"</div>"
    )


def build_pptx_slides_preview_html(file_path):
    if Presentation is None or file_path.suffix.lower() != ".pptx":
        return None

    presentation = Presentation(str(file_path))
    slide_width = emu_to_px(presentation.slide_width)
    slide_height = emu_to_px(presentation.slide_height)
    slides_markup = []
    nav_markup = []
    for index, slide in enumerate(presentation.slides, start=1):
        items = []
        collect_pptx_slide_items(
            slide,
            slide.shapes,
            {
                "origin_x": 0,
                "origin_y": 0,
                "child_origin_x": 0,
                "child_origin_y": 0,
                "scale_x": 1.0,
                "scale_y": 1.0,
            },
            items,
        )
        items.sort(key=lambda item: (item["top"], item["left"]))
        slide_title = detect_slide_title(slide, items) or f"第 {index} 页"
        slide_markup = "".join(render_pptx_slide_item(file_path, item) for item in items)
        slides_markup.append(
            f"""
            <section class="ppt-slide-panel {'active' if index == 1 else ''}" data-ppt-slide="{index}">
              <div class="ppt-slide-meta">
                <span class="ppt-slide-chip">Slide {index}</span>
                <strong>{html.escape(slide_title)}</strong>
              </div>
              <div class="ppt-slide-stage" style="aspect-ratio:{slide_width}/{slide_height};">
                <div class="ppt-slide-canvas" data-slide-width="{slide_width:.2f}" data-slide-height="{slide_height:.2f}">
                  {slide_markup}
                </div>
              </div>
            </section>
            """
        )
        nav_markup.append(
            f'<button class="ppt-nav-pill {"active" if index == 1 else ""}" type="button" data-ppt-jump="{index}">{index:02d}</button>'
        )

    title = html.escape(get_display_file_title(file_path))
    raw_url = build_upload_url(file_path.name)
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    :root {{
      color-scheme: light;
      --bg: #f8fafc;
      --surface: #ffffff;
      --text-main: #0f172a;
      --text-muted: #64748b;
      --border: #e2e8f0;
      --accent: #6366f1;
      --accent-soft: rgba(99, 102, 241, 0.12);
      --shadow: 0 14px 34px rgba(15, 23, 42, 0.08);
      --radius: 22px;
    }}
    * {{ box-sizing: border-box; }}
    html, body {{
      margin: 0;
      min-height: 100%;
      background: var(--bg);
      color: var(--text-main);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      overflow-x: hidden;
    }}
    body {{ padding: 18px; }}
    .ppt-preview-shell {{
      min-height: calc(100vh - 36px);
      display: flex;
      flex-direction: column;
      gap: 16px;
    }}
    .ppt-preview-topbar {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 18px 20px;
      border: 1px solid var(--border);
      border-radius: var(--radius);
      background: rgba(255, 255, 255, 0.95);
      box-shadow: var(--shadow);
    }}
    .ppt-preview-meta {{
      display: flex;
      flex-direction: column;
      gap: 8px;
      min-width: 0;
    }}
    .ppt-preview-badge {{
      display: inline-flex;
      align-self: flex-start;
      padding: 8px 12px;
      border-radius: 999px;
      background: var(--accent-soft);
      color: var(--accent);
      font-size: 13px;
      font-weight: 700;
      letter-spacing: 0.08em;
    }}
    .ppt-preview-title {{
      margin: 0;
      font-size: 30px;
      line-height: 1.2;
      word-break: break-word;
    }}
    .ppt-preview-hint {{
      margin: 0;
      font-size: 14px;
      color: var(--text-muted);
      line-height: 1.7;
    }}
    .ppt-preview-link {{
      display: inline-flex;
      align-items: center;
      justify-content: center;
      padding: 12px 18px;
      border-radius: 14px;
      background: var(--surface);
      border: 1px solid var(--border);
      color: var(--text-main);
      text-decoration: none;
      font-weight: 700;
      white-space: nowrap;
    }}
    .ppt-reader {{
      flex: 1;
      min-height: 0;
      border: 1px solid var(--border);
      border-radius: 28px;
      background: linear-gradient(180deg, rgba(255,255,255,0.96), rgba(244,247,251,0.98));
      box-shadow: var(--shadow);
      display: flex;
      flex-direction: column;
      overflow: hidden;
    }}
    .ppt-reader-head {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 18px 22px 12px;
      border-bottom: 1px solid rgba(226, 232, 240, 0.8);
    }}
    .ppt-reader-status {{
      font-size: 14px;
      color: var(--text-muted);
      font-weight: 600;
    }}
    .ppt-reader-actions {{
      display: flex;
      align-items: center;
      gap: 10px;
    }}
    .ppt-mini-btn {{
      border: 1px solid var(--border);
      background: var(--surface);
      color: var(--text-main);
      border-radius: 12px;
      padding: 10px 14px;
      font-weight: 700;
      cursor: pointer;
    }}
    .ppt-reader-body {{
      flex: 1;
      min-height: 0;
      padding: 12px 22px 18px;
      overflow: auto;
    }}
    .ppt-slide-panel {{
      display: none;
      gap: 12px;
    }}
    .ppt-slide-panel.active {{
      display: flex;
      flex-direction: column;
    }}
    .ppt-slide-meta {{
      display: flex;
      align-items: center;
      gap: 10px;
      flex-wrap: wrap;
    }}
    .ppt-slide-chip {{
      display: inline-flex;
      padding: 6px 10px;
      border-radius: 999px;
      background: rgba(15, 23, 42, 0.06);
      color: var(--text-muted);
      font-size: 12px;
      font-weight: 700;
      letter-spacing: 0.08em;
      text-transform: uppercase;
    }}
    .ppt-slide-stage {{
      position: relative;
      width: 100%;
      overflow: hidden;
      border-radius: 22px;
      background: #eef2ff;
      box-shadow: inset 0 0 0 1px rgba(99, 102, 241, 0.08);
    }}
    .ppt-slide-canvas {{
      position: absolute;
      top: 0;
      left: 0;
      background: #ffffff;
      transform-origin: top left;
      overflow: hidden;
      box-shadow: 0 12px 30px rgba(15, 23, 42, 0.08);
    }}
    .ppt-render-text {{
      overflow: hidden;
      white-space: normal;
      word-break: break-word;
      line-height: 1.32;
      border-radius: 16px;
    }}
    .ppt-render-image {{
      object-fit: contain;
      display: block;
    }}
    .ppt-render-box {{
      border-radius: 18px;
    }}
    .ppt-render-table-wrap {{
      overflow: hidden;
      border-radius: 16px;
      background: rgba(255, 255, 255, 0.96);
      border: 1px solid rgba(148, 163, 184, 0.24);
    }}
    .ppt-render-table {{
      width: 100%;
      height: 100%;
      border-collapse: collapse;
      table-layout: fixed;
      font-size: 16px;
    }}
    .ppt-render-table td {{
      border: 1px solid rgba(148, 163, 184, 0.24);
      padding: 10px 12px;
      vertical-align: top;
      word-break: break-word;
    }}
    .ppt-slide-nav {{
      display: flex;
      gap: 10px;
      padding: 0 22px 18px;
      overflow-x: auto;
      scrollbar-width: thin;
    }}
    .ppt-nav-pill {{
      flex: 0 0 auto;
      min-width: 52px;
      border: 1px solid var(--border);
      background: rgba(255,255,255,0.92);
      color: var(--text-muted);
      border-radius: 999px;
      padding: 10px 12px;
      font-weight: 700;
      cursor: pointer;
    }}
    .ppt-nav-pill.active {{
      background: var(--accent);
      color: #ffffff;
      border-color: transparent;
      box-shadow: 0 10px 18px rgba(99, 102, 241, 0.24);
    }}
    @media (max-width: 720px) {{
      body {{ padding: 12px; }}
      .ppt-preview-shell {{ min-height: calc(100vh - 24px); }}
      .ppt-preview-topbar {{
        flex-direction: column;
        align-items: stretch;
      }}
      .ppt-preview-link {{ width: 100%; }}
      .ppt-reader-head {{
        flex-direction: column;
        align-items: stretch;
      }}
      .ppt-reader-actions {{
        justify-content: space-between;
      }}
    }}
  </style>
</head>
<body>
  <main class="ppt-preview-shell">
    <header class="ppt-preview-topbar">
      <div class="ppt-preview-meta">
        <span class="ppt-preview-badge">PPT 学习视图</span>
        <h1 class="ppt-preview-title">{title}</h1>
        <p class="ppt-preview-hint">系统已将幻灯片转换为适合当前阅读区的学习视图，内容会按窗口宽度自动适配。</p>
      </div>
      <a class="ppt-preview-link" href="{raw_url}" target="_blank" rel="noreferrer">下载原文件</a>
    </header>
    <section class="ppt-reader">
      <div class="ppt-reader-head">
        <span class="ppt-reader-status" data-ppt-status>第 1 / {len(slides_markup)} 页</span>
        <div class="ppt-reader-actions">
          <button class="ppt-mini-btn" type="button" data-ppt-prev>上一页</button>
          <button class="ppt-mini-btn" type="button" data-ppt-next>下一页</button>
        </div>
      </div>
      <div class="ppt-reader-body">
        {''.join(slides_markup)}
      </div>
      <div class="ppt-slide-nav">
        {''.join(nav_markup)}
      </div>
    </section>
  </main>
  <script>
    (function () {{
      var panels = Array.prototype.slice.call(document.querySelectorAll('[data-ppt-slide]'));
      var navs = Array.prototype.slice.call(document.querySelectorAll('[data-ppt-jump]'));
      var statusNode = document.querySelector('[data-ppt-status]');
      var current = 0;

      function fitSlides() {{
        panels.forEach(function (panel) {{
          var stage = panel.querySelector('.ppt-slide-stage');
          var canvas = panel.querySelector('.ppt-slide-canvas');
          if (!stage || !canvas) return;
          var baseWidth = parseFloat(canvas.dataset.slideWidth || '0') || 1920;
          var baseHeight = parseFloat(canvas.dataset.slideHeight || '0') || 1080;
          var availableWidth = stage.clientWidth;
          var scale = availableWidth / baseWidth;
          canvas.style.width = baseWidth + 'px';
          canvas.style.height = baseHeight + 'px';
          canvas.style.transform = 'scale(' + scale + ')';
        }});
      }}

      function updateView() {{
        panels.forEach(function (panel, index) {{
          panel.classList.toggle('active', index === current);
        }});
        navs.forEach(function (button, index) {{
          button.classList.toggle('active', index === current);
        }});
        if (statusNode) {{
          statusNode.textContent = '第 ' + (current + 1) + ' / ' + panels.length + ' 页';
        }}
        fitSlides();
      }}

      document.querySelector('[data-ppt-prev]')?.addEventListener('click', function () {{
        current = (current - 1 + panels.length) % panels.length;
        updateView();
      }});
      document.querySelector('[data-ppt-next]')?.addEventListener('click', function () {{
        current = (current + 1) % panels.length;
        updateView();
      }});
      navs.forEach(function (button, index) {{
        button.addEventListener('click', function () {{
          current = index;
          updateView();
        }});
      }});
      window.addEventListener('resize', fitSlides);
      updateView();
    }})();
  </script>
</body>
</html>
"""


def normalize_quicklook_numeric_css(raw_html):
    property_pattern = (
        r"(?P<prop>\b(?:top|left|right|bottom|width|height|min-width|min-height|max-width|max-height|"
        r"margin(?:-(?:top|right|bottom|left))?|padding(?:-(?:top|right|bottom|left))?|"
        r"font-size|letter-spacing|border-width|border-radius)\s*:\s*)"
        r"(?P<value>-?\d+(?:\.\d+)?)"
        r"(?P<tail>\s*;)"
    )
    return re.sub(property_pattern, lambda m: f"{m.group('prop')}{m.group('value')}px{m.group('tail')}", raw_html)


def transform_quicklook_assets(raw_html):
    def replace_pdf_img(match):
        before = match.group("before") or ""
        src = match.group("src")
        after = match.group("after") or ""
        return (
            f'<embed class="ql-pdf-asset"{before}src="{src}#toolbar=0&navpanes=0&scrollbar=0&view=FitH"'
            f'{after} type="application/pdf" />'
        )

    return re.sub(
        r'<img(?P<before>[^>]*?)src="(?P<src>[^"]+\.pdf)"(?P<after>[^>]*)>',
        replace_pdf_img,
        raw_html,
        flags=re.IGNORECASE,
    )


def build_quicklook_runtime_assets(base_href):
    style = f"""
<base href="{html.escape(base_href)}" />
<style>
html, body {{
  margin: 0 !important;
  padding: 0 !important;
  background: #f8fafc !important;
  overflow-x: hidden !important;
}}
body {{
  padding: 18px !important;
}}
div.slide,
div.loading-slide {{
  margin: 0 !important;
  box-shadow: 0 16px 28px rgba(15, 23, 42, 0.08) !important;
  transform-origin: top left !important;
}}
.ql-slide-wrap {{
  position: relative;
  margin: 0 auto 18px;
  overflow: hidden;
}}
.ql-pdf-asset {{
  border: 0;
  background: transparent;
  display: block;
}}
img {{
  display: block;
}}
</style>
"""
    script = """
<script>
(function () {
  function fitSlides() {
    var viewportWidth = Math.max(document.documentElement.clientWidth || 0, window.innerWidth || 0);
    var availableWidth = Math.max(viewportWidth - 36, 320);
    var slides = document.querySelectorAll('div.slide, div.loading-slide');
    slides.forEach(function (slide) {
      var wrapper = slide.parentElement;
      if (!wrapper || !wrapper.classList.contains('ql-slide-wrap')) {
        wrapper = document.createElement('div');
        wrapper.className = 'ql-slide-wrap';
        slide.parentNode.insertBefore(wrapper, slide);
        wrapper.appendChild(slide);
      }
      if (!slide.dataset.baseWidth) {
        var computed = window.getComputedStyle(slide);
        slide.dataset.baseWidth = String(parseFloat(computed.width) || slide.scrollWidth || 1440);
        slide.dataset.baseHeight = String(parseFloat(computed.height) || slide.scrollHeight || 810);
      }
      var baseWidth = parseFloat(slide.dataset.baseWidth) || 1440;
      var baseHeight = parseFloat(slide.dataset.baseHeight) || 810;
      var scale = Math.min(availableWidth / baseWidth, 1);
      wrapper.style.width = (baseWidth * scale) + 'px';
      wrapper.style.height = (baseHeight * scale) + 'px';
      slide.style.width = baseWidth + 'px';
      slide.style.height = baseHeight + 'px';
      slide.style.transform = 'scale(' + scale + ')';
    });
  }
  window.addEventListener('load', fitSlides);
  window.addEventListener('resize', fitSlides);
})();
</script>
"""
    return style, script


def build_browser_ready_quicklook_html(file_path, preview_html_path):
    raw_html = preview_html_path.read_text(encoding="utf-8", errors="ignore")
    raw_html = normalize_quicklook_numeric_css(raw_html)
    raw_html = transform_quicklook_assets(raw_html)
    base_href = build_upload_url(f"{file_path.name}.qlpreview/")
    style, script = build_quicklook_runtime_assets(base_href)
    if "<head>" in raw_html:
        raw_html = raw_html.replace("<head>", f"<head>{style}", 1)
    else:
        raw_html = f"<head>{style}</head>{raw_html}"
    if "</body>" in raw_html:
        raw_html = raw_html.replace("</body>", f"{script}</body>", 1)
    else:
        raw_html = f"{raw_html}{script}"
    return raw_html


def get_preview_mode(file_path):
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    if suffix in {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".html", ".htm", ".docx", ".pptx"}:
        return "text"
    return "unsupported"


def get_preview_label(file_path):
    suffix = file_path.suffix.lower()
    labels = {
        ".pdf": "PDF 在线预览",
        ".pptx": "PPT 文字预览",
        ".docx": "Word 文字预览",
        ".txt": "文本预览",
        ".md": "Markdown 预览",
        ".csv": "表格文本预览",
        ".json": "JSON 文本预览",
    }
    return labels.get(suffix, "课件预览")


def build_courseware_preview_html(file_path):
    display_title = get_display_file_title(file_path)
    title = html.escape(display_title)
    raw_url = build_upload_url(file_path.name)
    preview_mode = get_preview_mode(file_path)
    preview_label = html.escape(get_preview_label(file_path))
    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        pptx_preview = build_pptx_slides_preview_html(file_path)
        if pptx_preview:
            return pptx_preview

    if suffix in QUICKLOOK_PREVIEW_SUFFIXES:
        quicklook_dir = ensure_quicklook_preview(file_path)
        quicklook_html = quicklook_dir / "Preview.html" if quicklook_dir else None
        if quicklook_html and quicklook_html.exists():
            preview_content = f"""
                <iframe
                  class="preview-document-frame preview-slides-frame"
                  src="{build_quicklook_viewer_url(file_path.name)}"
                  title="{title}"
                ></iframe>
            """
            preview_hint = "当前 PPT 已自动转换为浏览器可读的幻灯片预览。"
            preview_label = "PPT 幻灯片预览"
            return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    :root {{
      color-scheme: light;
      --bg: #f8fafc;
      --surface: #ffffff;
      --text-main: #0f172a;
      --text-muted: #64748b;
      --border: #e2e8f0;
      --accent: #6366f1;
      --shadow: 0 12px 32px rgba(15, 23, 42, 0.08);
      --radius: 22px;
    }}
    * {{
      box-sizing: border-box;
    }}
    html,
    body {{
      margin: 0;
      min-height: 100%;
      background: var(--bg);
      color: var(--text-main);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }}
    body {{
      padding: 20px;
    }}
    .preview-shell {{
      min-height: calc(100vh - 40px);
      display: flex;
      flex-direction: column;
      gap: 16px;
    }}
    .preview-topbar {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 18px 20px;
      border: 1px solid var(--border);
      border-radius: var(--radius);
      background: rgba(255, 255, 255, 0.92);
      box-shadow: var(--shadow);
    }}
    .preview-meta {{
      display: flex;
      flex-direction: column;
      gap: 8px;
      min-width: 0;
    }}
    .preview-badge {{
      display: inline-flex;
      align-self: flex-start;
      padding: 8px 12px;
      border-radius: 999px;
      background: rgba(99, 102, 241, 0.12);
      color: var(--accent);
      font-size: 13px;
      font-weight: 700;
      letter-spacing: 0.08em;
    }}
    .preview-title {{
      margin: 0;
      font-size: 28px;
      line-height: 1.2;
      word-break: break-word;
    }}
    .preview-hint {{
      margin: 0;
      color: var(--text-muted);
      font-size: 14px;
      line-height: 1.7;
    }}
    .preview-link {{
      display: inline-flex;
      align-items: center;
      justify-content: center;
      padding: 12px 18px;
      border-radius: 14px;
      background: var(--surface);
      border: 1px solid var(--border);
      color: var(--text-main);
      text-decoration: none;
      font-weight: 700;
      white-space: nowrap;
    }}
    .preview-body {{
      flex: 1;
      min-height: 0;
      border: 1px solid var(--border);
      border-radius: 26px;
      background: var(--surface);
      box-shadow: var(--shadow);
      overflow: hidden;
    }}
    .preview-document-frame {{
      width: 100%;
      height: 100%;
      min-height: calc(100vh - 188px);
      border: 0;
      background: #f1f5f9;
    }}
    .preview-slides-frame {{
      background: #d7dce5;
    }}
    @media (max-width: 720px) {{
      body {{
        padding: 12px;
      }}
      .preview-shell {{
        min-height: calc(100vh - 24px);
      }}
      .preview-topbar {{
        flex-direction: column;
        align-items: stretch;
      }}
      .preview-link {{
        width: 100%;
      }}
      .preview-body,
      .preview-document-frame {{
        min-height: calc(100vh - 216px);
      }}
    }}
  </style>
</head>
<body>
  <main class="preview-shell">
    <header class="preview-topbar">
      <div class="preview-meta">
        <span class="preview-badge">{html.escape(preview_label)}</span>
        <h1 class="preview-title">{title}</h1>
        <p class="preview-hint">{html.escape(preview_hint)}</p>
      </div>
      <a class="preview-link" href="{raw_url}" target="_blank" rel="noreferrer">下载原文件</a>
    </header>
    <section class="preview-body">
      {preview_content}
    </section>
  </main>
</body>
</html>
"""

    if preview_mode == "pdf":
        preview_content = f"""
            <iframe
              class="preview-document-frame"
              src="{raw_url}#toolbar=0&navpanes=0"
              title="{title}"
            ></iframe>
        """
        preview_hint = "当前课件为 PDF，已直接嵌入阅读区。"
    elif preview_mode == "image":
        preview_content = f"""
            <div class="preview-image-wrap">
              <img class="preview-image" src="{raw_url}" alt="{title}" />
            </div>
        """
        preview_hint = "当前课件为图片文件，可直接在线查看。"
    elif preview_mode == "text":
        extracted_text = extract_courseware_text(file_path)
        if extracted_text:
            preview_content = f"""
                <section class="preview-text-wrap">
                  <pre class="preview-text">{html.escape(extracted_text)}</pre>
                </section>
            """
            if file_path.suffix.lower() in QUICKLOOK_PREVIEW_SUFFIXES:
                preview_hint = "当前环境未生成幻灯片版式预览，已自动回退为文字预览。"
            elif file_path.suffix.lower() == ".docx":
                preview_hint = "已自动提取文档正文内容供在线阅读。"
            else:
                preview_hint = "已将课件正文内容整理为可读预览。"
        else:
            preview_content = """
                <div class="preview-empty-state">
                  <strong>暂时无法直接预览这份课件</strong>
                  <p>当前文件没有可提取的正文内容，建议下载原文件后使用本地软件打开。</p>
                </div>
            """
            preview_hint = "该课件暂不支持直接内嵌显示。"
    else:
        preview_content = """
            <div class="preview-empty-state">
              <strong>当前格式暂不支持在线预览</strong>
              <p>你可以先下载原文件，再使用本地应用打开。</p>
            </div>
        """
        preview_hint = "该格式在浏览器中通常会直接下载，系统已切换为安全预览页。"

    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    :root {{
      color-scheme: light;
      --bg: #f8fafc;
      --surface: #ffffff;
      --text-main: #0f172a;
      --text-muted: #64748b;
      --border: #e2e8f0;
      --accent: #6366f1;
      --shadow: 0 12px 32px rgba(15, 23, 42, 0.08);
      --radius: 22px;
    }}
    * {{
      box-sizing: border-box;
    }}
    html,
    body {{
      margin: 0;
      min-height: 100%;
      background: var(--bg);
      color: var(--text-main);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }}
    body {{
      padding: 20px;
    }}
    .preview-shell {{
      min-height: calc(100vh - 40px);
      display: flex;
      flex-direction: column;
      gap: 16px;
    }}
    .preview-topbar {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 18px 20px;
      border: 1px solid var(--border);
      border-radius: var(--radius);
      background: rgba(255, 255, 255, 0.92);
      box-shadow: var(--shadow);
    }}
    .preview-meta {{
      display: flex;
      flex-direction: column;
      gap: 8px;
      min-width: 0;
    }}
    .preview-badge {{
      display: inline-flex;
      align-self: flex-start;
      padding: 8px 12px;
      border-radius: 999px;
      background: rgba(99, 102, 241, 0.12);
      color: var(--accent);
      font-size: 13px;
      font-weight: 700;
      letter-spacing: 0.08em;
    }}
    .preview-title {{
      margin: 0;
      font-size: 28px;
      line-height: 1.2;
      word-break: break-word;
    }}
    .preview-hint {{
      margin: 0;
      color: var(--text-muted);
      font-size: 14px;
      line-height: 1.7;
    }}
    .preview-link {{
      display: inline-flex;
      align-items: center;
      justify-content: center;
      padding: 12px 18px;
      border-radius: 14px;
      background: var(--surface);
      border: 1px solid var(--border);
      color: var(--text-main);
      text-decoration: none;
      font-weight: 700;
      white-space: nowrap;
    }}
    .preview-body {{
      flex: 1;
      min-height: 0;
      border: 1px solid var(--border);
      border-radius: 26px;
      background: var(--surface);
      box-shadow: var(--shadow);
      overflow: hidden;
    }}
    .preview-document-frame {{
      width: 100%;
      height: 100%;
      min-height: calc(100vh - 188px);
      border: 0;
      background: #f1f5f9;
    }}
    .preview-image-wrap {{
      height: 100%;
      min-height: calc(100vh - 188px);
      display: flex;
      align-items: center;
      justify-content: center;
      padding: 24px;
      background: #f8fafc;
    }}
    .preview-image {{
      max-width: 100%;
      max-height: calc(100vh - 236px);
      object-fit: contain;
      border-radius: 18px;
      box-shadow: 0 10px 30px rgba(15, 23, 42, 0.12);
    }}
    .preview-text-wrap {{
      height: 100%;
      min-height: calc(100vh - 188px);
      padding: 28px 30px;
      overflow: auto;
      background:
        linear-gradient(180deg, rgba(99, 102, 241, 0.04), rgba(99, 102, 241, 0)),
        var(--surface);
    }}
    .preview-text {{
      margin: 0;
      white-space: pre-wrap;
      word-break: break-word;
      color: var(--text-main);
      font: 15px/1.8 -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }}
    .preview-empty-state {{
      min-height: calc(100vh - 188px);
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      gap: 12px;
      padding: 32px;
      text-align: center;
      color: var(--text-muted);
    }}
    .preview-empty-state strong {{
      color: var(--text-main);
      font-size: 18px;
    }}
    .preview-empty-state p {{
      margin: 0;
      max-width: 520px;
      line-height: 1.7;
    }}
    @media (max-width: 720px) {{
      body {{
        padding: 12px;
      }}
      .preview-shell {{
        min-height: calc(100vh - 24px);
      }}
      .preview-topbar {{
        flex-direction: column;
        align-items: stretch;
      }}
      .preview-link {{
        width: 100%;
      }}
      .preview-body,
      .preview-document-frame,
      .preview-image-wrap,
      .preview-text-wrap,
      .preview-empty-state {{
        min-height: calc(100vh - 216px);
      }}
    }}
  </style>
</head>
<body>
  <main class="preview-shell">
    <header class="preview-topbar">
      <div class="preview-meta">
        <span class="preview-badge">{preview_label}</span>
        <h1 class="preview-title">{title}</h1>
        <p class="preview-hint">{html.escape(preview_hint)}</p>
      </div>
      <a class="preview-link" href="{raw_url}" target="_blank" rel="noreferrer">下载原文件</a>
    </header>
    <section class="preview-body">
      {preview_content}
    </section>
  </main>
</body>
</html>
"""


def get_conn():
    conn = sqlite3.connect(DB_PATH, timeout=10)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    conn.execute("PRAGMA busy_timeout = 10000")
    return conn


def normalize_text_content(text):
    text = html.unescape(text or "")
    text = text.replace("\x00", " ")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_xml_text(xml_bytes):
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return ""
    texts = [node.text for node in root.iter() if node.text and node.text.strip()]
    return normalize_text_content("\n".join(texts))


def extract_pdf_text_fallback(file_path):
    data = file_path.read_bytes()
    chunks = []
    for match in re.finditer(rb"\(([^()]{2,400})\)", data):
        try:
            text = match.group(1).decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = match.group(1).decode("latin-1")
            except UnicodeDecodeError:
                continue
        if any(char.isalpha() or "\u4e00" <= char <= "\u9fff" for char in text):
            chunks.append(text)
    return normalize_text_content("\n".join(chunks))


def extract_pdf_text(file_path):
    if PdfReader is not None:
        try:
            reader = PdfReader(str(file_path))
            texts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    texts.append(page_text)
            parsed = normalize_text_content("\n\n".join(texts))
            if parsed:
                return parsed
        except Exception:
            pass
    return extract_pdf_text_fallback(file_path)


def extract_courseware_text(file_path):
    if not file_path.exists():
        return ""

    suffix = file_path.suffix.lower()
    text_like_suffixes = {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".html", ".htm"}
    if suffix in text_like_suffixes:
        raw = file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix in {".html", ".htm"}:
            raw = re.sub(r"<[^>]+>", " ", raw)
        return normalize_text_content(raw)

    if suffix == ".docx":
        try:
            with zipfile.ZipFile(file_path) as archive:
                return extract_xml_text(archive.read("word/document.xml"))
        except (KeyError, zipfile.BadZipFile, OSError):
            return ""

    if suffix == ".pptx":
        try:
            texts = []
            with zipfile.ZipFile(file_path) as archive:
                slide_names = sorted(
                    name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                )
                for name in slide_names:
                    slide_text = extract_xml_text(archive.read(name))
                    if slide_text:
                        texts.append(slide_text)
            return normalize_text_content("\n\n".join(texts))
        except (zipfile.BadZipFile, OSError):
            return ""

    if suffix == ".pdf":
        return extract_pdf_text(file_path)

    return ""


def crop_ai_context(text, limit=MAX_AI_CONTEXT_CHARS):
    text = normalize_text_content(text)
    if len(text) <= limit:
        return text
    return f"{text[:limit].rstrip()}\n\n[以下内容已截断]"


def stringify_model_content(content):
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
            elif isinstance(item, dict):
                text = item.get("text")
                if not text and isinstance(item.get("content"), str):
                    text = item["content"]
                if text:
                    parts.append(str(text))
        return "\n".join(part.strip() for part in parts if str(part).strip()).strip()
    return str(content).strip() if content is not None else ""


def extract_bigmodel_answer(message):
    if not isinstance(message, dict):
        return ""
    content = stringify_model_content(message.get("content"))
    if content:
        return content
    return stringify_model_content(message.get("reasoning_content"))


def parse_bigmodel_error(payload):
    if isinstance(payload, dict):
        error = payload.get("error")
        if isinstance(error, dict):
            return error.get("message") or error.get("code")
        if isinstance(error, str):
            return error
        return payload.get("message")
    return None


def call_bigmodel_chat(messages):
    api_key = os.environ.get("BIGMODEL_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("服务器未配置 BIGMODEL_API_KEY，暂时无法使用 AI 问答。")

    payload = {
        "model": BIGMODEL_MODEL,
        "messages": messages,
        "temperature": 0.3,
        "max_tokens": 2048,
        "thinking": {"type": "disabled"},
    }
    request = urllib.request.Request(
        BIGMODEL_API_URL,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=90) as response:
            data = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as error:
        detail = error.read().decode("utf-8", errors="ignore")
        try:
            payload = json.loads(detail)
        except json.JSONDecodeError:
            payload = None
        message = parse_bigmodel_error(payload) or f"模型服务调用失败（HTTP {error.code}）。"
        raise RuntimeError(message) from error
    except urllib.error.URLError as error:
        raise RuntimeError("模型服务暂时不可用，请稍后再试。") from error

    try:
        content = extract_bigmodel_answer(data["choices"][0]["message"])
    except (KeyError, IndexError, TypeError) as error:
        raise RuntimeError("模型返回格式异常，暂时无法生成回答。") from error

    if not content:
        raise RuntimeError("模型未返回有效内容，请稍后重试。")
    return content


def publish_user_updates(*user_ids):
    global SYNC_CURSOR

    listeners = {int(user_id) for user_id in user_ids if user_id}
    if not listeners:
        return

    with SYNC_CONDITION:
        SYNC_CURSOR += 1
        for user_id in listeners:
            USER_SYNC_CURSORS[user_id] = SYNC_CURSOR
        SYNC_CONDITION.notify_all()


def get_user_sync_cursor(user_id):
    return USER_SYNC_CURSORS.get(user_id, 0)


def wait_for_user_update(user_id, cursor, timeout=20):
    deadline = time.monotonic() + timeout
    with SYNC_CONDITION:
        while USER_SYNC_CURSORS.get(user_id, 0) <= cursor:
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                break
            SYNC_CONDITION.wait(remaining)
        return USER_SYNC_CURSORS.get(user_id, 0)


def hash_password(password, salt=None):
    salt = salt or secrets.token_hex(16)
    digest = hashlib.pbkdf2_hmac(
        "sha256",
        password.encode("utf-8"),
        salt.encode("utf-8"),
        100_000,
    ).hex()
    return salt, digest


def verify_password(password, salt, password_hash):
    _, digest = hash_password(password, salt)
    return secrets.compare_digest(digest, password_hash)


def row_to_user(row):
    return {
        "id": row["id"],
        "username": row["username"],
        "display_name": row["display_name"],
        "role": row["role"],
        "student_number": row["student_number"],
        "created_at": row["created_at"],
    }


def ensure_column(conn, table_name, column_name, definition):
    columns = {row["name"] for row in conn.execute(f"PRAGMA table_info({table_name})").fetchall()}
    if column_name not in columns:
        conn.execute(f"ALTER TABLE {table_name} ADD COLUMN {column_name} {definition}")


def recreate_users_table_with_admin_role(conn):
    columns = {row["name"] for row in conn.execute("PRAGMA table_info(users)").fetchall()}
    student_number_select = "student_number" if "student_number" in columns else "NULL AS student_number"

    conn.commit()
    conn.execute("PRAGMA foreign_keys = OFF")
    conn.execute(
        """
        CREATE TABLE users_new (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            display_name TEXT NOT NULL,
            role TEXT NOT NULL CHECK(role IN ('teacher', 'student', 'admin')),
            student_number TEXT,
            salt TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
        """
    )
    conn.execute(
        f"""
        INSERT INTO users_new (id, username, display_name, role, student_number, salt, password_hash, created_at)
        SELECT id, username, display_name, role, {student_number_select}, salt, password_hash, created_at
        FROM users
        """
    )
    conn.execute("DROP TABLE users")
    conn.execute("ALTER TABLE users_new RENAME TO users")
    conn.commit()
    conn.execute("PRAGMA foreign_keys = ON")


def ensure_users_table_supports_admin(conn):
    schema_row = conn.execute(
        """
        SELECT sql
        FROM sqlite_master
        WHERE type = 'table' AND name = 'users'
        """
    ).fetchone()
    schema_sql = (schema_row["sql"] or "") if schema_row else ""
    if "'admin'" not in schema_sql:
        recreate_users_table_with_admin_role(conn)


def normalize_user_pair(user_a, user_b):
    return (user_a, user_b) if user_a < user_b else (user_b, user_a)


def find_thread_id(conn, user_a, user_b):
    user_one, user_two = normalize_user_pair(user_a, user_b)
    row = conn.execute(
        """
        SELECT id
        FROM conversation_threads
        WHERE user_one_id = ? AND user_two_id = ?
        """,
        (user_one, user_two),
    ).fetchone()
    return row["id"] if row else None


def seed_classes(conn):
    class_rows = conn.execute("SELECT id, teacher_id FROM classes").fetchall()
    for row in class_rows:
        conn.execute(
            "INSERT OR IGNORE INTO class_members (class_id, user_id, joined_at) VALUES (?, ?, ?)",
            (row["id"], row["teacher_id"], now_iso()),
        )
    conn.commit()


def backfill_class_links(conn):
    default_class = conn.execute("SELECT id FROM classes ORDER BY id LIMIT 1").fetchone()
    if not default_class:
        return

    default_class_id = default_class["id"]
    conn.execute("UPDATE coursewares SET class_id = ? WHERE class_id IS NULL", (default_class_id,))
    conn.execute("UPDATE discussions SET class_id = ? WHERE class_id IS NULL", (default_class_id,))
    conn.commit()


def get_or_create_thread(conn, user_a, user_b, visible_for_a=1, visible_for_b=1):
    thread_id = find_thread_id(conn, user_a, user_b)
    if not thread_id:
        user_one, user_two = normalize_user_pair(user_a, user_b)
        cursor = conn.execute(
            """
            INSERT INTO conversation_threads (user_one_id, user_two_id, created_at, last_message_at)
            VALUES (?, ?, ?, ?)
            """,
            (user_one, user_two, now_iso(), now_iso()),
        )
        thread_id = cursor.lastrowid

    conn.execute(
        """
        INSERT OR IGNORE INTO conversation_members (thread_id, user_id, visible, joined_at)
        VALUES (?, ?, ?, ?)
        """,
        (thread_id, user_a, visible_for_a, now_iso()),
    )
    conn.execute(
        """
        INSERT OR IGNORE INTO conversation_members (thread_id, user_id, visible, joined_at)
        VALUES (?, ?, ?, ?)
        """,
        (thread_id, user_b, visible_for_b, now_iso()),
    )
    return thread_id


def migrate_message_threads(conn):
    rows = conn.execute(
        """
        SELECT id, sender_id, receiver_id, created_at
        FROM messages
        WHERE thread_id IS NULL
        ORDER BY id ASC
        """
    ).fetchall()

    for row in rows:
        thread_id = get_or_create_thread(conn, row["sender_id"], row["receiver_id"])
        conn.execute(
            "UPDATE messages SET thread_id = ? WHERE id = ?",
            (thread_id, row["id"]),
        )
        conn.execute(
            "UPDATE conversation_threads SET last_message_at = ? WHERE id = ?",
            (row["created_at"], thread_id),
        )
    conn.commit()


def build_class_summary(conn, class_id):
    row = conn.execute(
        """
        SELECT
            c.id,
            c.name,
            c.description,
            c.teacher_id,
            c.created_at,
            u.display_name AS teacher_name,
            (
                SELECT COUNT(*)
                FROM class_members cm
                JOIN users member_user ON member_user.id = cm.user_id
                WHERE cm.class_id = c.id AND member_user.role = 'student'
            ) AS student_count,
            (
                SELECT COUNT(*)
                FROM class_join_requests r
                WHERE r.class_id = c.id AND r.status = 'pending'
            ) AS pending_request_count
        FROM classes c
        JOIN users u ON u.id = c.teacher_id
        WHERE c.id = ?
        """,
        (class_id,),
    ).fetchone()
    return dict(row) if row else None


def user_can_access_class(conn, user, class_id):
    if not class_id:
        return False
    if user["role"] == "admin":
        return True
    if user["role"] == "teacher":
        row = conn.execute(
            "SELECT id FROM classes WHERE id = ? AND teacher_id = ?",
            (class_id, user["id"]),
        ).fetchone()
        if row:
            return True
    member = conn.execute(
        "SELECT 1 FROM class_members WHERE class_id = ? AND user_id = ?",
        (class_id, user["id"]),
    ).fetchone()
    return bool(member)


def teacher_can_manage_class(conn, teacher_id, class_id):
    row = conn.execute(
        "SELECT id FROM classes WHERE id = ? AND teacher_id = ?",
        (class_id, teacher_id),
    ).fetchone()
    return bool(row)


def user_can_manage_class(conn, user, class_id):
    return user["role"] == "admin" or teacher_can_manage_class(conn, user["id"], class_id)


def users_share_class(conn, user_a, user_b):
    row = conn.execute(
        """
        SELECT 1
        FROM class_members cm1
        JOIN class_members cm2 ON cm2.class_id = cm1.class_id
        WHERE cm1.user_id = ? AND cm2.user_id = ?
        LIMIT 1
        """,
        (user_a, user_b),
    ).fetchone()
    return bool(row)


def build_class_scope_clause(user, class_id, column_name):
    if class_id is not None:
        return f"{column_name} = ?", [class_id]
    if user["role"] == "admin":
        return "1 = 1", []
    if user["role"] == "teacher":
        return f"{column_name} IN (SELECT id FROM classes WHERE teacher_id = ?)", [user["id"]]
    return f"{column_name} IN (SELECT class_id FROM class_members WHERE user_id = ?)", [user["id"]]


def get_join_request(conn, class_id, student_id):
    return conn.execute(
        """
        SELECT id, class_id, student_id, status, requested_at, reviewed_at, reviewed_by
        FROM class_join_requests
        WHERE class_id = ? AND student_id = ?
        """,
        (class_id, student_id),
    ).fetchone()


def init_db():
    DATA_DIR.mkdir(exist_ok=True)
    UPLOAD_DIR.mkdir(exist_ok=True)

    conn = get_conn()
    conn.execute("PRAGMA journal_mode = WAL")
    conn.executescript(
        """
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL UNIQUE,
            display_name TEXT NOT NULL,
            role TEXT NOT NULL CHECK(role IN ('teacher', 'student', 'admin')),
            student_number TEXT,
            salt TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS sessions (
            token TEXT PRIMARY KEY,
            user_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(user_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS classes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            description TEXT,
            teacher_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(teacher_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS class_members (
            class_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            joined_at TEXT NOT NULL,
            PRIMARY KEY (class_id, user_id),
            FOREIGN KEY(class_id) REFERENCES classes(id),
            FOREIGN KEY(user_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS class_join_requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            class_id INTEGER NOT NULL,
            student_id INTEGER NOT NULL,
            status TEXT NOT NULL CHECK(status IN ('pending', 'approved', 'rejected')),
            requested_at TEXT NOT NULL,
            reviewed_at TEXT,
            reviewed_by INTEGER,
            UNIQUE(class_id, student_id),
            FOREIGN KEY(class_id) REFERENCES classes(id),
            FOREIGN KEY(student_id) REFERENCES users(id),
            FOREIGN KEY(reviewed_by) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS coursewares (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            course_name TEXT NOT NULL,
            description TEXT,
            original_file_name TEXT NOT NULL,
            stored_file_name TEXT NOT NULL,
            uploaded_by INTEGER NOT NULL,
            uploaded_at TEXT NOT NULL,
            FOREIGN KEY(uploaded_by) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS evaluations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            courseware_id INTEGER NOT NULL,
            student_id INTEGER NOT NULL,
            helpfulness INTEGER NOT NULL,
            usability INTEGER NOT NULL,
            suitability INTEGER NOT NULL DEFAULT 3,
            practicality INTEGER NOT NULL DEFAULT 3,
            suggestion TEXT,
            created_at TEXT NOT NULL,
            UNIQUE(courseware_id, student_id),
            FOREIGN KEY(courseware_id) REFERENCES coursewares(id),
            FOREIGN KEY(student_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS discussions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            body TEXT NOT NULL,
            author_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(author_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS discussion_replies (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            discussion_id INTEGER NOT NULL,
            body TEXT NOT NULL,
            author_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(discussion_id) REFERENCES discussions(id),
            FOREIGN KEY(author_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS conversation_threads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_one_id INTEGER NOT NULL,
            user_two_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            last_message_at TEXT NOT NULL,
            UNIQUE(user_one_id, user_two_id),
            FOREIGN KEY(user_one_id) REFERENCES users(id),
            FOREIGN KEY(user_two_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS conversation_members (
            thread_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            visible INTEGER NOT NULL DEFAULT 1,
            joined_at TEXT NOT NULL,
            PRIMARY KEY (thread_id, user_id),
            FOREIGN KEY(thread_id) REFERENCES conversation_threads(id),
            FOREIGN KEY(user_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            sender_id INTEGER NOT NULL,
            receiver_id INTEGER NOT NULL,
            body TEXT NOT NULL,
            is_read INTEGER NOT NULL DEFAULT 0,
            created_at TEXT NOT NULL,
            FOREIGN KEY(sender_id) REFERENCES users(id),
            FOREIGN KEY(receiver_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS ai_chat_messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            courseware_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            role TEXT NOT NULL CHECK(role IN ('user', 'assistant')),
            content TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(courseware_id) REFERENCES coursewares(id),
            FOREIGN KEY(user_id) REFERENCES users(id)
        );

        CREATE TABLE IF NOT EXISTS rag_chat_messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            class_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            role TEXT NOT NULL CHECK(role IN ('user', 'assistant')),
            content TEXT NOT NULL,
            sources TEXT NOT NULL DEFAULT '[]',
            created_at TEXT NOT NULL,
            FOREIGN KEY(class_id) REFERENCES classes(id),
            FOREIGN KEY(user_id) REFERENCES users(id)
        );
        """
    )
    ensure_users_table_supports_admin(conn)
    ensure_column(conn, "users", "student_number", "TEXT")
    ensure_column(conn, "coursewares", "class_id", "INTEGER")
    ensure_column(conn, "discussions", "class_id", "INTEGER")
    ensure_column(conn, "messages", "thread_id", "INTEGER")
    ensure_column(conn, "evaluations", "suitability", "INTEGER")
    ensure_column(conn, "evaluations", "practicality", "INTEGER")
    conn.execute(
        """
        CREATE UNIQUE INDEX IF NOT EXISTS idx_users_student_number
        ON users(student_number)
        WHERE student_number IS NOT NULL AND student_number != ''
        """
    )
    conn.execute(
        """
        CREATE INDEX IF NOT EXISTS idx_ai_chat_messages_courseware_user
        ON ai_chat_messages(courseware_id, user_id, id)
        """
    )
    conn.execute(
        """
        CREATE INDEX IF NOT EXISTS idx_rag_chat_messages_class_user
        ON rag_chat_messages(class_id, user_id, id)
        """
    )
    conn.execute(
        """
        UPDATE users
        SET student_number = CASE username
            WHEN 'student01' THEN '20260001'
            WHEN 'student02' THEN '20260002'
            ELSE student_number
        END
        WHERE role = 'student' AND (student_number IS NULL OR student_number = '')
          AND username IN ('student01', 'student02')
        """
    )
    conn.execute(
        """
        UPDATE evaluations
        SET suitability = COALESCE(suitability, usability),
            practicality = COALESCE(practicality, helpfulness)
        WHERE suitability IS NULL OR practicality IS NULL
        """
    )
    conn.commit()

    seed_demo_users(conn)
    ensure_admin_user(conn)
    seed_classes(conn)
    seed_courseware(conn)
    seed_discussion(conn)
    seed_messages(conn)
    backfill_class_links(conn)
    migrate_message_threads(conn)
    conn.close()


def seed_demo_users(conn):
    existing = conn.execute("SELECT COUNT(*) AS count FROM users").fetchone()["count"]
    if existing:
        return

    demo_users = [
        ("teacher01", "王老师", "teacher", None, "Teacher@123"),
        ("student01", "李同学", "student", "20260001", "Student@123"),
        ("student02", "陈同学", "student", "20260002", "Student@123"),
    ]
    for username, display_name, role, student_number, password in demo_users:
        salt, password_hash = hash_password(password)
        conn.execute(
            """
            INSERT INTO users (username, display_name, role, student_number, salt, password_hash, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (username, display_name, role, student_number, salt, password_hash, now_iso()),
        )
    conn.commit()


def ensure_admin_user(conn):
    salt, password_hash = hash_password(ADMIN_PASSWORD)
    existing_admin = conn.execute(
        """
        SELECT id
        FROM users
        WHERE role = 'admin'
        ORDER BY id ASC
        LIMIT 1
        """
    ).fetchone()

    if existing_admin:
        conn.execute(
            """
            UPDATE users
            SET username = ?, display_name = ?, student_number = NULL, salt = ?, password_hash = ?
            WHERE id = ?
            """,
            (ADMIN_USERNAME, ADMIN_DISPLAY_NAME, salt, password_hash, existing_admin["id"]),
        )
        conn.execute(
            "DELETE FROM sessions WHERE user_id != ? AND user_id IN (SELECT id FROM users WHERE role = 'admin')",
            (existing_admin["id"],),
        )
        conn.execute(
            "DELETE FROM users WHERE role = 'admin' AND id != ?",
            (existing_admin["id"],),
        )
        conn.commit()
        return

    conn.execute(
        """
        INSERT INTO users (username, display_name, role, student_number, salt, password_hash, created_at)
        VALUES (?, ?, 'admin', NULL, ?, ?, ?)
        """,
        (ADMIN_USERNAME, ADMIN_DISPLAY_NAME, salt, password_hash, now_iso()),
    )
    conn.commit()


def seed_courseware(conn):
    existing = conn.execute("SELECT COUNT(*) AS count FROM coursewares").fetchone()["count"]
    if existing:
        return

    teacher = conn.execute(
        "SELECT id FROM users WHERE role = 'teacher' ORDER BY id LIMIT 1"
    ).fetchone()
    if not teacher:
        return
    class_row = conn.execute(
        "SELECT id FROM classes WHERE teacher_id = ? ORDER BY id LIMIT 1",
        (teacher["id"],),
    ).fetchone()
    if not class_row:
        return
    class_id = class_row["id"]

    demo_file = UPLOAD_DIR / "demo_courseware.txt"
    if not demo_file.exists():
        demo_file.write_text(
            "AI助教系统演示课件\n\n"
            "1. 系统支持教师上传课件并统一管理。\n"
            "2. 学生可以围绕课件内容进行学习和讨论。\n"
            "3. AI问答界面当前为前端演示版，便于后续接入真实模型。\n",
            encoding="utf-8",
        )

    conn.execute(
        """
        INSERT INTO coursewares (
            title, course_name, description, original_file_name, stored_file_name, uploaded_by, uploaded_at, class_id
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            "AI助教系统导论",
            "软件工程课程设计",
            "用于演示课件上传、浏览和前端问答入口的示例课件。",
            "demo_courseware.txt",
            "demo_courseware.txt",
            teacher["id"],
            now_iso(),
            class_id,
        ),
    )
    conn.commit()


def seed_discussion(conn):
    existing = conn.execute("SELECT COUNT(*) AS count FROM discussions").fetchone()["count"]
    if existing:
        return

    teacher = conn.execute(
        "SELECT id FROM users WHERE role = 'teacher' ORDER BY id LIMIT 1"
    ).fetchone()
    student = conn.execute(
        "SELECT id FROM users WHERE role = 'student' ORDER BY id LIMIT 1"
    ).fetchone()
    if not teacher or not student:
        return
    class_row = conn.execute(
        "SELECT id FROM classes WHERE teacher_id = ? ORDER BY id LIMIT 1",
        (teacher["id"],),
    ).fetchone()
    if not class_row:
        return
    class_id = class_row["id"]
    is_member = conn.execute(
        "SELECT 1 FROM class_members WHERE class_id = ? AND user_id = ?",
        (class_id, student["id"]),
    ).fetchone()
    if not is_member:
        return

    cursor = conn.execute(
        """
        INSERT INTO discussions (title, body, author_id, created_at, class_id)
        VALUES (?, ?, ?, ?, ?)
        """,
        (
            "第一次使用 AI 助教系统时应该先看什么？",
            "建议大家先阅读示例课件，再体验 AI 问答、问卷和讨论区，这样能更快理解系统闭环。",
            student["id"],
            now_iso(),
            class_id,
        ),
    )
    discussion_id = cursor.lastrowid
    conn.execute(
        """
        INSERT INTO discussion_replies (discussion_id, body, author_id, created_at)
        VALUES (?, ?, ?, ?)
        """,
        (
            discussion_id,
            "可以先从课件详情页进入，里面已经预留了 AI 问答入口，后续也方便扩展真实模型能力。",
            teacher["id"],
            now_iso(),
        ),
    )
    conn.commit()


def seed_messages(conn):
    existing = conn.execute("SELECT COUNT(*) AS count FROM messages").fetchone()["count"]
    if existing:
        return

    teacher = conn.execute(
        "SELECT id FROM users WHERE role = 'teacher' ORDER BY id LIMIT 1"
    ).fetchone()
    student = conn.execute(
        "SELECT id FROM users WHERE role = 'student' ORDER BY id LIMIT 1"
    ).fetchone()
    if not teacher or not student:
        return
    if not users_share_class(conn, teacher["id"], student["id"]):
        return

    conn.executemany(
        """
        INSERT INTO messages (sender_id, receiver_id, body, is_read, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        [
            (student["id"], teacher["id"], "老师好，我已经试用了系统，后续 AI 回答会接入真实模型吗？", 1, now_iso()),
            (teacher["id"], student["id"], "会的，现在先把教学流程和页面原型跑通，后面再接模型和向量检索。", 0, now_iso()),
        ],
    )
    conn.commit()


def delete_courseware_record(conn, courseware_id):
    row = conn.execute(
        "SELECT stored_file_name, class_id FROM coursewares WHERE id = ?",
        (courseware_id,),
    ).fetchone()
    if not row:
        return {"stored_file_names": [], "affected_user_ids": []}

    affected_user_ids = [
        member["user_id"]
        for member in conn.execute(
            "SELECT user_id FROM class_members WHERE class_id = ?",
            (row["class_id"],),
        ).fetchall()
    ]
    conn.execute("DELETE FROM ai_chat_messages WHERE courseware_id = ?", (courseware_id,))
    conn.execute("DELETE FROM evaluations WHERE courseware_id = ?", (courseware_id,))
    conn.execute("DELETE FROM coursewares WHERE id = ?", (courseware_id,))
    return {"stored_file_names": [row["stored_file_name"]], "affected_user_ids": affected_user_ids}


def delete_class_record(conn, class_id):
    class_row = conn.execute(
        "SELECT id FROM classes WHERE id = ?",
        (class_id,),
    ).fetchone()
    if not class_row:
        return {"stored_file_names": [], "affected_user_ids": []}

    affected_user_ids = {
        row["user_id"]
        for row in conn.execute(
            "SELECT user_id FROM class_members WHERE class_id = ?",
            (class_id,),
        ).fetchall()
    }
    affected_user_ids.update(
        row["student_id"]
        for row in conn.execute(
            "SELECT student_id FROM class_join_requests WHERE class_id = ?",
            (class_id,),
        ).fetchall()
    )

    stored_file_names = [
        row["stored_file_name"]
        for row in conn.execute(
            "SELECT stored_file_name FROM coursewares WHERE class_id = ?",
            (class_id,),
        ).fetchall()
    ]
    conn.execute(
        "DELETE FROM ai_chat_messages WHERE courseware_id IN (SELECT id FROM coursewares WHERE class_id = ?)",
        (class_id,),
    )
    conn.execute(
        "DELETE FROM evaluations WHERE courseware_id IN (SELECT id FROM coursewares WHERE class_id = ?)",
        (class_id,),
    )
    conn.execute("DELETE FROM coursewares WHERE class_id = ?", (class_id,))
    conn.execute(
        "DELETE FROM discussion_replies WHERE discussion_id IN (SELECT id FROM discussions WHERE class_id = ?)",
        (class_id,),
    )
    conn.execute("DELETE FROM discussions WHERE class_id = ?", (class_id,))
    conn.execute("DELETE FROM class_join_requests WHERE class_id = ?", (class_id,))
    conn.execute("DELETE FROM class_members WHERE class_id = ?", (class_id,))
    conn.execute("DELETE FROM classes WHERE id = ?", (class_id,))
    return {"stored_file_names": stored_file_names, "affected_user_ids": list(affected_user_ids)}


def delete_user_conversations(conn, user_id):
    thread_ids = [
        row["id"]
        for row in conn.execute(
            """
            SELECT id
            FROM conversation_threads
            WHERE user_one_id = ? OR user_two_id = ?
            """,
            (user_id, user_id),
        ).fetchall()
    ]
    affected_user_ids = {
        row["other_id"]
        for row in conn.execute(
            """
            SELECT CASE
                WHEN user_one_id = ? THEN user_two_id
                ELSE user_one_id
            END AS other_id
            FROM conversation_threads
            WHERE user_one_id = ? OR user_two_id = ?
            """,
            (user_id, user_id, user_id),
        ).fetchall()
    }

    conn.execute("DELETE FROM messages WHERE sender_id = ? OR receiver_id = ?", (user_id, user_id))
    if thread_ids:
        placeholders = ", ".join("?" for _ in thread_ids)
        conn.execute(
            f"DELETE FROM conversation_members WHERE thread_id IN ({placeholders})",
            thread_ids,
        )
        conn.execute(
            f"DELETE FROM conversation_threads WHERE id IN ({placeholders})",
            thread_ids,
        )
    else:
        conn.execute("DELETE FROM conversation_members WHERE user_id = ?", (user_id,))
    return list(affected_user_ids)


def delete_user_record(conn, user_id):
    user_row = conn.execute(
        "SELECT id, role FROM users WHERE id = ?",
        (user_id,),
    ).fetchone()
    if not user_row:
        return {"stored_file_names": [], "affected_user_ids": []}

    stored_file_names = []
    affected_user_ids = set()

    class_ids = [
        row["id"]
        for row in conn.execute(
            "SELECT id FROM classes WHERE teacher_id = ?",
            (user_id,),
        ).fetchall()
    ]
    for class_id in class_ids:
        result = delete_class_record(conn, class_id)
        stored_file_names.extend(result["stored_file_names"])
        affected_user_ids.update(result["affected_user_ids"])

    affected_user_ids.update(
        row["user_id"]
        for row in conn.execute(
            """
            SELECT DISTINCT cm2.user_id
            FROM class_members cm1
            JOIN class_members cm2 ON cm2.class_id = cm1.class_id
            WHERE cm1.user_id = ? AND cm2.user_id != ?
            """,
            (user_id, user_id),
        ).fetchall()
    )

    uploaded_courseware_ids = [
        row["id"]
        for row in conn.execute(
            "SELECT id FROM coursewares WHERE uploaded_by = ?",
            (user_id,),
        ).fetchall()
    ]
    for courseware_id in uploaded_courseware_ids:
        result = delete_courseware_record(conn, courseware_id)
        stored_file_names.extend(result["stored_file_names"])
        affected_user_ids.update(result["affected_user_ids"])

    discussion_ids = [
        row["id"]
        for row in conn.execute(
            "SELECT id FROM discussions WHERE author_id = ?",
            (user_id,),
        ).fetchall()
    ]
    if discussion_ids:
        placeholders = ", ".join("?" for _ in discussion_ids)
        conn.execute(
            f"DELETE FROM discussion_replies WHERE discussion_id IN ({placeholders})",
            discussion_ids,
        )
        conn.execute(
            f"DELETE FROM discussions WHERE id IN ({placeholders})",
            discussion_ids,
        )
    conn.execute("DELETE FROM discussion_replies WHERE author_id = ?", (user_id,))
    conn.execute("DELETE FROM ai_chat_messages WHERE user_id = ?", (user_id,))
    conn.execute("DELETE FROM evaluations WHERE student_id = ?", (user_id,))
    conn.execute("DELETE FROM class_join_requests WHERE student_id = ? OR reviewed_by = ?", (user_id, user_id))
    conn.execute("DELETE FROM class_members WHERE user_id = ?", (user_id,))
    conn.execute("DELETE FROM sessions WHERE user_id = ?", (user_id,))
    affected_user_ids.update(delete_user_conversations(conn, user_id))
    conn.execute("DELETE FROM users WHERE id = ?", (user_id,))

    return {
        "stored_file_names": stored_file_names,
        "affected_user_ids": [user for user in affected_user_ids if user != user_id],
    }


class AppHandler(BaseHTTPRequestHandler):
    server_version = "AITutorServer/1.0"

    def do_GET(self):
        parsed = urlparse(self.path)
        path = parsed.path

        if path.startswith("/api/"):
            self.handle_api_get(path)
            return

        if path.startswith("/preview/"):
            self.serve_courseware_preview(path.removeprefix("/preview/"))
            return

        if path.startswith("/preview-quicklook/"):
            self.serve_quicklook_browser_preview(path.removeprefix("/preview-quicklook/"))
            return

        if path == "/preview-media":
            self.serve_ppt_media_asset()
            return

        if path.startswith("/uploads/"):
            self.serve_file(UPLOAD_DIR, path.removeprefix("/uploads/"))
            return

        self.serve_static(path)

    def do_POST(self):
        parsed = urlparse(self.path)
        path = parsed.path

        if path.startswith("/api/"):
            self.handle_api_post(path)
            return

        self.send_error(HTTPStatus.NOT_FOUND, "Not found")

    def do_PUT(self):
        parsed = urlparse(self.path)
        path = parsed.path

        if path.startswith("/api/"):
            self.handle_api_put(path)
            return

        self.send_error(HTTPStatus.NOT_FOUND, "Not found")

    def do_DELETE(self):
        parsed = urlparse(self.path)
        path = parsed.path

        if path.startswith("/api/"):
            self.handle_api_delete(path)
            return

        self.send_error(HTTPStatus.NOT_FOUND, "Not found")

    def do_OPTIONS(self):
        self.send_response(HTTPStatus.NO_CONTENT)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Headers", "Content-Type, Authorization")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, PUT, DELETE, OPTIONS")
        self.end_headers()

    def log_message(self, fmt, *args):
        return

    def get_current_user(self):
        auth_header = self.headers.get("Authorization", "")
        if not auth_header.startswith("Bearer "):
            return None

        token = auth_header.split(" ", 1)[1].strip()
        if not token:
            return None

        conn = get_conn()
        row = conn.execute(
            """
            SELECT u.*
            FROM sessions s
            JOIN users u ON u.id = s.user_id
            WHERE s.token = ?
            """,
            (token,),
        ).fetchone()
        conn.close()
        return row_to_user(row) if row else None

    def require_user(self):
        user = self.get_current_user()
        if not user:
            self.send_json({"error": "请先登录。"}, status=HTTPStatus.UNAUTHORIZED)
            return None
        return user

    def require_role(self, role):
        roles = tuple(role) if isinstance(role, (list, tuple, set)) else (role,)
        user = self.require_user()
        if not user:
            return None
        if user["role"] not in roles:
            self.send_json({"error": "当前角色没有权限执行该操作。"}, status=HTTPStatus.FORBIDDEN)
            return None
        return user

    def parse_json_body(self):
        length = int(self.headers.get("Content-Length", "0") or 0)
        if length <= 0:
            return {}
        raw = self.rfile.read(length)
        if not raw:
            return {}
        return json.loads(raw.decode("utf-8"))

    def get_query_params(self):
        parsed = urlparse(self.path)
        return {key: values[0] for key, values in parse_qs(parsed.query).items()}

    def get_query_int(self, key):
        value = self.get_query_params().get(key)
        if value in (None, ""):
            return None
        try:
            return int(value)
        except ValueError:
            return None

    def parse_multipart_form(self):
        content_type = self.headers.get("Content-Type", "")
        length = int(self.headers.get("Content-Length", "0") or 0)
        body = self.rfile.read(length)
        message = BytesParser(policy=default).parsebytes(
            f"Content-Type: {content_type}\r\nMIME-Version: 1.0\r\n\r\n".encode("utf-8") + body
        )

        fields = {}
        for part in message.iter_parts():
            name = part.get_param("name", header="Content-Disposition")
            if not name:
                continue

            filename = part.get_filename()
            payload = part.get_payload(decode=True) or b""
            fields[name] = {
                "filename": filename,
                "data": payload,
                "value": payload.decode("utf-8", errors="ignore"),
            }
        return fields

    def send_json(self, payload, status=HTTPStatus.OK):
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(body)

    def send_html(self, html_text, status=HTTPStatus.OK):
        body = html_text.encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.end_headers()
        self.wfile.write(body)

    def send_bytes(self, body, content_type="application/octet-stream", status=HTTPStatus.OK):
        self.send_response(status)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.end_headers()
        self.wfile.write(body)

    def require_class_access(self, conn, user, class_id):
        if class_id is None:
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return False
        if not user_can_access_class(conn, user, class_id):
            self.send_json({"error": "当前用户无权访问该班级。"}, status=HTTPStatus.FORBIDDEN)
            return False
        return True

    def serve_static(self, path):
        relative = "index.html" if path in ("", "/", "/index.html") else path.lstrip("/")
        file_path = (STATIC_DIR / relative).resolve()

        if not str(file_path).startswith(str(STATIC_DIR.resolve())) or not file_path.exists():
            file_path = STATIC_DIR / "index.html"

        self.send_file(file_path)

    def serve_courseware_preview(self, relative_path):
        relative_path = unquote(relative_path).lstrip("/")
        stored_name = relative_path.split("/", 1)[0]
        file_path = (UPLOAD_DIR / stored_name).resolve()
        if not str(file_path).startswith(str(UPLOAD_DIR.resolve())) or not file_path.exists():
            self.send_error(HTTPStatus.NOT_FOUND, "File not found")
            return
        if file_path.suffix.lower() == ".pdf":
            display_name = f"{get_display_file_title(file_path)}.pdf"
            self.send_file(file_path, download_name=display_name, inline=True)
            return
        if file_path.suffix.lower() in QUICKLOOK_PREVIEW_SUFFIXES:
            native_preview = ensure_native_pdf_preview(file_path)
            if native_preview and native_preview.exists():
                display_name = f"{get_display_file_title(file_path)}.pdf"
                self.send_file(native_preview, download_name=display_name, inline=True)
                return
        self.send_html(build_courseware_preview_html(file_path))

    def serve_quicklook_browser_preview(self, relative_path):
        relative_path = unquote(relative_path).lstrip("/")
        file_path = (UPLOAD_DIR / relative_path).resolve()
        if not str(file_path).startswith(str(UPLOAD_DIR.resolve())) or not file_path.exists():
            self.send_error(HTTPStatus.NOT_FOUND, "File not found")
            return
        preview_dir = ensure_quicklook_preview(file_path)
        preview_html = preview_dir / "Preview.html" if preview_dir else None
        if not preview_html or not preview_html.exists():
            self.send_error(HTTPStatus.NOT_FOUND, "Preview not found")
            return
        self.send_html(build_browser_ready_quicklook_html(file_path, preview_html))

    def serve_ppt_media_asset(self):
        params = self.get_query_params()
        file_name = (params.get("file") or "").strip()
        asset_name = os.path.basename((params.get("asset") or "").strip())
        if not file_name or not asset_name:
            self.send_error(HTTPStatus.BAD_REQUEST, "Invalid preview asset request")
            return

        file_path = (UPLOAD_DIR / unquote(file_name)).resolve()
        if not str(file_path).startswith(str(UPLOAD_DIR.resolve())) or not file_path.exists():
            self.send_error(HTTPStatus.NOT_FOUND, "File not found")
            return

        media_path = f"ppt/media/{asset_name}"
        try:
            with zipfile.ZipFile(file_path) as archive:
                body = archive.read(media_path)
        except (OSError, zipfile.BadZipFile, KeyError):
            self.send_error(HTTPStatus.NOT_FOUND, "Preview asset not found")
            return

        content_type, _ = mimetypes.guess_type(asset_name)
        self.send_bytes(body, content_type or "application/octet-stream")

    def serve_file(self, base_dir, relative_path):
        relative_path = unquote(relative_path).lstrip("/")
        file_path = (base_dir / relative_path).resolve()
        if not str(file_path).startswith(str(base_dir.resolve())) or not file_path.exists():
            self.send_error(HTTPStatus.NOT_FOUND, "File not found")
            return
        self.send_file(file_path)

    def send_file(self, file_path, content_type=None, download_name=None, inline=False):
        guessed_type, _ = mimetypes.guess_type(file_path.name)
        content_type = content_type or guessed_type
        content_type = content_type or "application/octet-stream"
        data = file_path.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        if download_name:
            download_path = Path(download_name)
            safe_stem = re.sub(r"[^A-Za-z0-9_-]+", "_", download_path.stem).strip("._-") or "preview"
            safe_suffix = download_path.suffix if download_path.suffix else ""
            fallback_name = f"{safe_stem}{safe_suffix}"
            disposition_type = "inline" if inline else "attachment"
            self.send_header(
                "Content-Disposition",
                f"{disposition_type}; filename=\"{fallback_name}\"; filename*=UTF-8''{quote(download_name)}",
            )
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.end_headers()
        self.wfile.write(data)

    def handle_api_get(self, path):
        if path == "/api/me":
            return self.api_me()
        if path == "/api/classes":
            return self.api_classes()
        if path == "/api/classes/available":
            return self.api_available_classes()
        if path.startswith("/api/classes/") and path.endswith("/members"):
            return self.api_class_members(path)
        if path == "/api/dashboard":
            return self.api_dashboard()
        if path == "/api/users":
            return self.api_users()
        if path == "/api/coursewares":
            return self.api_coursewares()
        if path.startswith("/api/coursewares/"):
            return self.api_courseware_detail(path)
        if path == "/api/evaluations":
            return self.api_evaluations()
        if path == "/api/ai/messages":
            return self.api_ai_messages()
        if path == "/api/rag/status":
            return self.api_rag_status()
        if path == "/api/rag/messages":
            return self.api_rag_messages()
        if path == "/api/discussions":
            return self.api_discussions()
        if path == "/api/messages/contacts":
            return self.api_message_contacts()
        if path == "/api/messages/conversations":
            return self.api_conversations()
        if path == "/api/messages/events":
            return self.api_message_events()
        if path.startswith("/api/messages/thread/"):
            return self.api_message_thread(path)

        self.send_json({"error": "接口不存在。"}, status=HTTPStatus.NOT_FOUND)

    def handle_api_post(self, path):
        if path == "/api/auth/login":
            return self.api_login()
        if path == "/api/auth/register":
            return self.api_register()
        if path == "/api/auth/logout":
            return self.api_logout()
        if path == "/api/classes":
            return self.api_create_class()
        if path == "/api/classes/join":
            return self.api_join_class()
        if path.startswith("/api/classes/requests/") and path.endswith("/approve"):
            return self.api_review_class_request(path, "approved")
        if path.startswith("/api/classes/requests/") and path.endswith("/reject"):
            return self.api_review_class_request(path, "rejected")
        if path.startswith("/api/classes/") and path.endswith("/members"):
            return self.api_add_class_member(path)
        if path == "/api/users":
            return self.api_create_user()
        if path == "/api/coursewares":
            return self.api_create_courseware()
        if path == "/api/evaluations":
            return self.api_create_evaluation()
        if path == "/api/ai/messages":
            return self.api_ai_ask()
        if path == "/api/rag/index":
            return self.api_rag_index()
        if path == "/api/rag/ask":
            return self.api_rag_ask()
        if path == "/api/discussions":
            return self.api_create_discussion()
        if path.startswith("/api/discussions/") and path.endswith("/replies"):
            return self.api_create_discussion_reply(path)
        if path == "/api/messages/conversations":
            return self.api_create_conversation()
        if path == "/api/messages":
            return self.api_send_message()

        self.send_json({"error": "接口不存在。"}, status=HTTPStatus.NOT_FOUND)

    def handle_api_put(self, path):
        if path.startswith("/api/users/"):
            return self.api_update_user(path)
        if path.startswith("/api/classes/"):
            return self.api_update_class(path)
        if path.startswith("/api/coursewares/"):
            return self.api_update_courseware(path)
        self.send_json({"error": "接口不存在。"}, status=HTTPStatus.NOT_FOUND)

    def handle_api_delete(self, path):
        if path.startswith("/api/classes/") and "/members/" in path:
            return self.api_remove_class_member(path)
        if path.startswith("/api/classes/"):
            return self.api_delete_class(path)
        if path.startswith("/api/users/"):
            return self.api_delete_user(path)
        if path.startswith("/api/coursewares/"):
            return self.api_delete_courseware(path)
        if path == "/api/ai/messages":
            return self.api_clear_ai_messages()
        if path == "/api/rag/messages":
            return self.api_rag_clear_messages()
        if path.startswith("/api/messages/conversations/"):
            return self.api_delete_conversation(path)
        self.send_json({"error": "接口不存在。"}, status=HTTPStatus.NOT_FOUND)

    def api_me(self):
        user = self.require_user()
        if not user:
            return
        self.send_json({"user": user})

    def api_classes(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        if user["role"] == "admin":
            rows = conn.execute(
                """
                SELECT id
                FROM classes
                ORDER BY created_at DESC, id DESC
                """
            ).fetchall()
        elif user["role"] == "teacher":
            rows = conn.execute(
                """
                SELECT id
                FROM classes
                WHERE teacher_id = ?
                ORDER BY created_at DESC, id DESC
                """,
                (user["id"],),
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT c.id
                FROM classes c
                JOIN class_members cm ON cm.class_id = c.id
                WHERE cm.user_id = ?
                ORDER BY c.created_at DESC, c.id DESC
                """,
                (user["id"],),
            ).fetchall()

        classes = []
        for row in rows:
            item = build_class_summary(conn, row["id"])
            if item:
                item["is_owner"] = user["role"] == "admin" or item["teacher_id"] == user["id"]
                classes.append(item)
        conn.close()
        self.send_json({"classes": classes})

    def api_available_classes(self):
        user = self.require_role("student")
        if not user:
            return

        conn = get_conn()
        rows = conn.execute(
            """
            SELECT
                c.id,
                r.id AS request_id,
                r.status AS join_request_status,
                r.requested_at,
                r.reviewed_at
            FROM classes c
            LEFT JOIN class_join_requests r
              ON r.class_id = c.id AND r.student_id = ? AND r.status != 'approved'
            WHERE c.id NOT IN (
                SELECT class_id
                FROM class_members
                WHERE user_id = ?
            )
            ORDER BY c.created_at DESC, c.id DESC
            """,
            (user["id"], user["id"]),
        ).fetchall()

        classes = []
        for row in rows:
            item = build_class_summary(conn, row["id"])
            if item:
                item["join_request_id"] = row["request_id"]
                item["join_request_status"] = row["join_request_status"]
                item["join_requested_at"] = row["requested_at"]
                item["join_reviewed_at"] = row["reviewed_at"]
                classes.append(item)
        conn.close()
        self.send_json({"classes": classes})

    def api_create_class(self):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        data = self.parse_json_body()
        name = (data.get("name") or "").strip()
        description = (data.get("description") or "").strip()
        teacher_id = data.get("teacher_id")
        if not name:
            self.send_json({"error": "请输入班级名称。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if user["role"] == "admin":
            try:
                teacher_id = int(teacher_id)
            except (TypeError, ValueError):
                conn.close()
                self.send_json({"error": "请选择授课教师。"}, status=HTTPStatus.BAD_REQUEST)
                return
            teacher = conn.execute(
                "SELECT id FROM users WHERE id = ? AND role = 'teacher'",
                (teacher_id,),
            ).fetchone()
            if not teacher:
                conn.close()
                self.send_json({"error": "授课教师不存在。"}, status=HTTPStatus.BAD_REQUEST)
                return
        else:
            teacher_id = user["id"]
        cursor = conn.execute(
            """
            INSERT INTO classes (name, description, teacher_id, created_at)
            VALUES (?, ?, ?, ?)
            """,
            (name, description, teacher_id, now_iso()),
        )
        class_id = cursor.lastrowid
        conn.execute(
            "INSERT OR IGNORE INTO class_members (class_id, user_id, joined_at) VALUES (?, ?, ?)",
            (class_id, teacher_id, now_iso()),
        )
        conn.commit()
        class_info = build_class_summary(conn, class_id)
        conn.close()
        publish_user_updates(user["id"], teacher_id)
        self.send_json({"classroom": class_info}, status=HTTPStatus.CREATED)

    def api_update_class(self, path):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        try:
            class_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "班级编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if not user_can_manage_class(conn, user, class_id):
            conn.close()
            self.send_json({"error": "当前教师无权修改该班级。"}, status=HTTPStatus.FORBIDDEN)
            return

        data = self.parse_json_body()
        name = (data.get("name") or "").strip()
        description = (data.get("description") or "").strip()
        teacher_id = data.get("teacher_id")
        if not name:
            conn.close()
            self.send_json({"error": "请输入班级名称。"}, status=HTTPStatus.BAD_REQUEST)
            return

        current_class = conn.execute(
            "SELECT teacher_id FROM classes WHERE id = ?",
            (class_id,),
        ).fetchone()
        if not current_class:
            conn.close()
            self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
            return

        new_teacher_id = current_class["teacher_id"]
        if user["role"] == "admin":
            try:
                new_teacher_id = int(teacher_id)
            except (TypeError, ValueError):
                conn.close()
                self.send_json({"error": "请选择授课教师。"}, status=HTTPStatus.BAD_REQUEST)
                return
            teacher = conn.execute(
                "SELECT id FROM users WHERE id = ? AND role = 'teacher'",
                (new_teacher_id,),
            ).fetchone()
            if not teacher:
                conn.close()
                self.send_json({"error": "授课教师不存在。"}, status=HTTPStatus.BAD_REQUEST)
                return

        conn.execute(
            """
            UPDATE classes
            SET name = ?, description = ?, teacher_id = ?
            WHERE id = ?
            """,
            (name, description, new_teacher_id, class_id),
        )
        conn.execute(
            "INSERT OR IGNORE INTO class_members (class_id, user_id, joined_at) VALUES (?, ?, ?)",
            (class_id, new_teacher_id, now_iso()),
        )
        if current_class["teacher_id"] != new_teacher_id:
            conn.execute(
                "DELETE FROM class_members WHERE class_id = ? AND user_id = ?",
                (class_id, current_class["teacher_id"]),
            )
        conn.commit()
        class_info = build_class_summary(conn, class_id)
        member_ids = [row["user_id"] for row in conn.execute("SELECT user_id FROM class_members WHERE class_id = ?", (class_id,))]
        conn.close()
        publish_user_updates(*member_ids, current_class["teacher_id"], new_teacher_id)
        self.send_json({"classroom": class_info, "message": "班级信息已更新。"})

    def api_delete_class(self, path):
        user = self.require_role("admin")
        if not user:
            return

        try:
            class_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "班级编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        class_row = conn.execute(
            "SELECT id FROM classes WHERE id = ?",
            (class_id,),
        ).fetchone()
        if not class_row:
            conn.close()
            self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
            return

        result = delete_class_record(conn, class_id)
        conn.commit()
        conn.close()

        for file_name in result["stored_file_names"]:
            delete_courseware_assets(file_name)

        publish_user_updates(*result["affected_user_ids"])
        self.send_json({"message": "班级已删除。"})

    def api_join_class(self):
        user = self.require_role("student")
        if not user:
            return

        data = self.parse_json_body()
        class_id = data.get("class_id")
        try:
            class_id = int(class_id)
        except (TypeError, ValueError):
            self.send_json({"error": "班级参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        classroom = conn.execute("SELECT id, teacher_id FROM classes WHERE id = ?", (class_id,)).fetchone()
        if not classroom:
            conn.close()
            self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
            return

        member = conn.execute(
            "SELECT 1 FROM class_members WHERE class_id = ? AND user_id = ?",
            (class_id, user["id"]),
        ).fetchone()
        if member:
            conn.close()
            self.send_json({"error": "你已经在该班级中。"}, status=HTTPStatus.BAD_REQUEST)
            return

        existing = get_join_request(conn, class_id, user["id"])
        if existing and existing["status"] == "pending":
            conn.close()
            self.send_json({"error": "你已提交申请，请等待教师审核。"}, status=HTTPStatus.BAD_REQUEST)
            return

        if existing:
            conn.execute(
                """
                UPDATE class_join_requests
                SET status = 'pending', requested_at = ?, reviewed_at = NULL, reviewed_by = NULL
                WHERE id = ?
                """,
                (now_iso(), existing["id"]),
            )
        else:
            conn.execute(
                """
                INSERT INTO class_join_requests (class_id, student_id, status, requested_at)
                VALUES (?, ?, 'pending', ?)
                """,
                (class_id, user["id"], now_iso()),
            )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"], classroom["teacher_id"])
        self.send_json({"message": "申请已提交，等待教师审核。"}, status=HTTPStatus.CREATED)

    def api_class_members(self, path):
        user = self.require_user()
        if not user:
            return

        try:
            class_id = int(path.split("/")[-2])
        except ValueError:
            self.send_json({"error": "班级编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        classroom = build_class_summary(conn, class_id)
        members = [
            dict(row)
            for row in conn.execute(
                """
                SELECT u.id, u.display_name, u.username, u.role, u.student_number, cm.joined_at
                FROM class_members cm
                JOIN users u ON u.id = cm.user_id
                WHERE cm.class_id = ?
                ORDER BY CASE u.role WHEN 'teacher' THEN 0 ELSE 1 END, u.display_name ASC
                """,
                (class_id,),
            ).fetchall()
        ]

        payload = {"classroom": classroom, "members": members}
        if user_can_manage_class(conn, user, class_id):
            available_students = [
                dict(row)
                for row in conn.execute(
                    """
                    SELECT id, display_name, username, role, student_number
                    FROM users
                    WHERE role = 'student'
                      AND id NOT IN (
                          SELECT user_id
                          FROM class_members
                          WHERE class_id = ?
                      )
                      AND id NOT IN (
                          SELECT student_id
                          FROM class_join_requests
                          WHERE class_id = ? AND status = 'pending'
                      )
                    ORDER BY display_name ASC
                    """,
                    (class_id, class_id),
                ).fetchall()
            ]
            pending_requests = [
                dict(row)
                for row in conn.execute(
                    """
                    SELECT
                        r.id,
                        r.student_id,
                        r.requested_at,
                        u.display_name,
                        u.username,
                        u.student_number
                    FROM class_join_requests r
                    JOIN users u ON u.id = r.student_id
                    WHERE r.class_id = ? AND r.status = 'pending'
                      AND r.student_id NOT IN (
                          SELECT user_id
                          FROM class_members
                          WHERE class_id = ?
                      )
                    ORDER BY r.requested_at ASC, r.id ASC
                    """,
                    (class_id, class_id),
                ).fetchall()
            ]
            payload["available_students"] = available_students
            payload["pending_requests"] = pending_requests

        conn.close()
        self.send_json(payload)

    def api_review_class_request(self, path, decision):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        try:
            request_id = int(path.split("/")[-2])
        except ValueError:
            self.send_json({"error": "申请编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        request_row = conn.execute(
            """
            SELECT r.id, r.class_id, r.student_id, r.status, c.teacher_id
            FROM class_join_requests r
            JOIN classes c ON c.id = r.class_id
            WHERE r.id = ?
            """,
            (request_id,),
        ).fetchone()
        if not request_row:
            conn.close()
            self.send_json({"error": "申请记录不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if user["role"] != "admin" and request_row["teacher_id"] != user["id"]:
            conn.close()
            self.send_json({"error": "当前教师无权审核该申请。"}, status=HTTPStatus.FORBIDDEN)
            return
        if request_row["status"] != "pending":
            conn.close()
            self.send_json({"error": "该申请已处理。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn.execute(
            """
            UPDATE class_join_requests
            SET status = ?, reviewed_at = ?, reviewed_by = ?
            WHERE id = ?
            """,
            (decision, now_iso(), user["id"], request_id),
        )
        if decision == "approved":
            conn.execute(
                "INSERT OR IGNORE INTO class_members (class_id, user_id, joined_at) VALUES (?, ?, ?)",
                (request_row["class_id"], request_row["student_id"], now_iso()),
            )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"], request_row["student_id"])
        message = "已通过加入申请。" if decision == "approved" else "已拒绝加入申请。"
        self.send_json({"message": message})

    def api_add_class_member(self, path):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        try:
            class_id = int(path.split("/")[-2])
        except ValueError:
            self.send_json({"error": "班级编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if not user_can_manage_class(conn, user, class_id):
            conn.close()
            self.send_json({"error": "当前教师无权管理该班级。"}, status=HTTPStatus.FORBIDDEN)
            return

        data = self.parse_json_body()
        student_id = data.get("student_id")
        try:
            student_id = int(student_id)
        except (TypeError, ValueError):
            conn.close()
            self.send_json({"error": "学生参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        student = conn.execute(
            "SELECT id, role FROM users WHERE id = ?",
            (student_id,),
        ).fetchone()
        if not student or student["role"] != "student":
            conn.close()
            self.send_json({"error": "仅可添加学生到班级。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn.execute(
            "INSERT OR IGNORE INTO class_members (class_id, user_id, joined_at) VALUES (?, ?, ?)",
            (class_id, student_id, now_iso()),
        )
        conn.execute(
            """
            UPDATE class_join_requests
            SET status = 'approved', reviewed_at = ?, reviewed_by = ?
            WHERE class_id = ? AND student_id = ?
            """,
            (now_iso(), user["id"], class_id, student_id),
        )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"], student_id)
        self.send_json({"message": "学生已加入班级。"}, status=HTTPStatus.CREATED)

    def api_remove_class_member(self, path):
        user = self.require_user()
        if not user:
            return

        parts = path.split("/")
        try:
            class_id = int(parts[-3])
            member_id = int(parts[-1])
        except ValueError:
            self.send_json({"error": "参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if user["id"] != member_id and not user_can_manage_class(conn, user, class_id):
            conn.close()
            self.send_json({"error": "当前用户无权移除该成员。"}, status=HTTPStatus.FORBIDDEN)
            return

        classroom = conn.execute(
            "SELECT teacher_id FROM classes WHERE id = ?",
            (class_id,),
        ).fetchone()
        if not classroom:
            conn.close()
            self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
            return

        if classroom["teacher_id"] == member_id:
            conn.close()
            self.send_json({"error": "班级教师不能从班级中移除。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn.execute(
            "DELETE FROM class_members WHERE class_id = ? AND user_id = ?",
            (class_id, member_id),
        )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"], member_id)
        self.send_json({"message": "成员已移除。"})

    def api_login(self):
        data = self.parse_json_body()
        username = (data.get("username") or "").strip()
        password = data.get("password") or ""
        if not username or not password:
            self.send_json({"error": "请输入用户名和密码。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        row = conn.execute("SELECT * FROM users WHERE username = ?", (username,)).fetchone()
        if not row or not verify_password(password, row["salt"], row["password_hash"]):
            conn.close()
            self.send_json({"error": "用户名或密码错误。"}, status=HTTPStatus.UNAUTHORIZED)
            return

        token = secrets.token_hex(24)
        conn.execute(
            "INSERT INTO sessions (token, user_id, created_at) VALUES (?, ?, ?)",
            (token, row["id"], now_iso()),
        )
        conn.commit()
        conn.close()
        self.send_json({"token": token, "user": row_to_user(row)})

    def api_register(self):
        data = self.parse_json_body()
        username = (data.get("username") or "").strip()
        display_name = (data.get("display_name") or "").strip()
        password = data.get("password") or ""
        role = (data.get("role") or "").strip()
        student_number = (data.get("student_number") or "").strip()

        if not username or not display_name or not password or role not in {"teacher", "student"}:
            self.send_json({"error": "请完整填写注册信息。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student" and not student_number:
            self.send_json({"error": "学生注册时请填写学号。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        exists = conn.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
        if exists:
            conn.close()
            self.send_json({"error": "用户名已存在。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student":
            duplicate_student_number = conn.execute(
                "SELECT id FROM users WHERE student_number = ?",
                (student_number,),
            ).fetchone()
            if duplicate_student_number:
                conn.close()
                self.send_json({"error": "学号已存在。"}, status=HTTPStatus.BAD_REQUEST)
                return

        salt, password_hash = hash_password(password)
        cursor = conn.execute(
            """
            INSERT INTO users (username, display_name, role, student_number, salt, password_hash, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                username,
                display_name,
                role,
                student_number if role == "student" else None,
                salt,
                password_hash,
                now_iso(),
            ),
        )
        user_id = cursor.lastrowid
        token = secrets.token_hex(24)
        conn.execute(
            "INSERT INTO sessions (token, user_id, created_at) VALUES (?, ?, ?)",
            (token, user_id, now_iso()),
        )
        user = conn.execute("SELECT * FROM users WHERE id = ?", (user_id,)).fetchone()
        conn.commit()
        conn.close()
        self.send_json({"token": token, "user": row_to_user(user)}, status=HTTPStatus.CREATED)

    def api_logout(self):
        user = self.require_user()
        if not user:
            return

        auth_header = self.headers.get("Authorization", "")
        token = auth_header.split(" ", 1)[1].strip()
        conn = get_conn()
        conn.execute("DELETE FROM sessions WHERE token = ?", (token,))
        conn.commit()
        conn.close()
        self.send_json({"message": "已退出登录。"})

    def api_dashboard(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        if user["role"] == "admin":
            teacher_count = conn.execute(
                "SELECT COUNT(*) AS count FROM users WHERE role = 'teacher'"
            ).fetchone()["count"]
            student_count = conn.execute(
                "SELECT COUNT(*) AS count FROM users WHERE role = 'student'"
            ).fetchone()["count"]
            class_count = conn.execute(
                "SELECT COUNT(*) AS count FROM classes"
            ).fetchone()["count"]
            courseware_count = conn.execute(
                "SELECT COUNT(*) AS count FROM coursewares"
            ).fetchone()["count"]
            pending_requests = conn.execute(
                "SELECT COUNT(*) AS count FROM class_join_requests WHERE status = 'pending'"
            ).fetchone()["count"]
            recent_classes = [
                dict(row)
                for row in conn.execute(
                    """
                    SELECT c.id, c.name, c.created_at, u.display_name AS teacher_name
                    FROM classes c
                    JOIN users u ON u.id = c.teacher_id
                    ORDER BY c.created_at DESC, c.id DESC
                    LIMIT 5
                    """
                ).fetchall()
            ]
            recent_coursewares = [
                dict(row)
                for row in conn.execute(
                    """
                    SELECT c.id, c.title, c.course_name, c.uploaded_at, cls.name AS class_name, u.display_name AS teacher_name
                    FROM coursewares c
                    JOIN classes cls ON cls.id = c.class_id
                    JOIN users u ON u.id = cls.teacher_id
                    ORDER BY c.uploaded_at DESC, c.id DESC
                    LIMIT 5
                    """
                ).fetchall()
            ]
            conn.close()
            self.send_json(
                {
                    "stats": {
                        "teachers": teacher_count,
                        "students": student_count,
                        "classes": class_count,
                        "coursewares": courseware_count,
                        "pending_requests": pending_requests,
                    },
                    "recent_classes": recent_classes,
                    "recent_coursewares": recent_coursewares,
                }
            )
            return

        class_id = self.get_query_int("class_id")
        if class_id is not None and not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        courseware_clause, courseware_params = build_class_scope_clause(user, class_id, "class_id")
        courseware_count = conn.execute(
            f"SELECT COUNT(*) AS count FROM coursewares WHERE {courseware_clause}",
            courseware_params,
        ).fetchone()["count"]
        discussion_count = conn.execute(
            f"SELECT COUNT(*) AS count FROM discussions WHERE {courseware_clause}",
            courseware_params,
        ).fetchone()["count"]
        recent_coursewares = [
            dict(row)
            for row in conn.execute(
                f"""
                SELECT c.id, c.title, c.course_name, c.uploaded_at, u.display_name AS teacher_name
                FROM coursewares c
                JOIN users u ON u.id = c.uploaded_by
                WHERE {courseware_clause}
                ORDER BY c.uploaded_at DESC
                LIMIT 5
                """,
                courseware_params,
            ).fetchall()
        ]

        if user["role"] == "teacher":
            evaluation_count = conn.execute(
                f"""
                SELECT COUNT(*) AS count
                FROM evaluations e
                JOIN coursewares c ON c.id = e.courseware_id
                WHERE {courseware_clause}
                """,
                courseware_params,
            ).fetchone()["count"]
            unread_messages = conn.execute(
                "SELECT COUNT(*) AS count FROM messages WHERE receiver_id = ? AND is_read = 0",
                (user["id"],),
            ).fetchone()["count"]
            payload = {
                "stats": {
                    "coursewares": courseware_count,
                    "evaluations": evaluation_count,
                    "discussions": discussion_count,
                    "unread_messages": unread_messages,
                },
                "recent_coursewares": recent_coursewares,
            }
        else:
            completed_surveys = conn.execute(
                f"""
                SELECT COUNT(*) AS count
                FROM evaluations e
                JOIN coursewares c ON c.id = e.courseware_id
                WHERE e.student_id = ? AND {courseware_clause}
                """,
                [user["id"], *courseware_params],
            ).fetchone()["count"]
            unread_messages = conn.execute(
                "SELECT COUNT(*) AS count FROM messages WHERE receiver_id = ? AND is_read = 0",
                (user["id"],),
            ).fetchone()["count"]
            payload = {
                "stats": {
                    "coursewares": courseware_count,
                    "completed_surveys": completed_surveys,
                    "discussions": discussion_count,
                    "unread_messages": unread_messages,
                },
                "recent_coursewares": recent_coursewares,
            }

        conn.close()
        self.send_json(payload)

    def api_users(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        if user["role"] == "admin":
            rows = conn.execute(
                """
                SELECT id, username, display_name, role, student_number, created_at
                FROM users
                ORDER BY CASE role WHEN 'admin' THEN 0 WHEN 'teacher' THEN 1 ELSE 2 END, created_at DESC, id DESC
                """
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT id, username, display_name, role, student_number, created_at
                FROM users
                WHERE id != ?
                ORDER BY CASE role WHEN 'teacher' THEN 0 ELSE 1 END, display_name ASC
                """,
                (user["id"],),
            ).fetchall()
        conn.close()
        self.send_json({"users": [dict(row) for row in rows]})

    def api_create_user(self):
        user = self.require_role("admin")
        if not user:
            return

        data = self.parse_json_body()
        username = (data.get("username") or "").strip()
        display_name = (data.get("display_name") or "").strip()
        password = data.get("password") or ""
        role = (data.get("role") or "").strip()
        student_number = (data.get("student_number") or "").strip()

        if not username or not display_name or not password or role not in {"teacher", "student"}:
            self.send_json({"error": "请完整填写用户信息。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student" and not student_number:
            self.send_json({"error": "创建学生时请填写学号。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        existing = conn.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
        if existing:
            conn.close()
            self.send_json({"error": "用户名已存在。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student":
            duplicate_student_number = conn.execute(
                "SELECT id FROM users WHERE student_number = ?",
                (student_number,),
            ).fetchone()
            if duplicate_student_number:
                conn.close()
                self.send_json({"error": "学号已存在。"}, status=HTTPStatus.BAD_REQUEST)
                return

        salt, password_hash = hash_password(password)
        cursor = conn.execute(
            """
            INSERT INTO users (username, display_name, role, student_number, salt, password_hash, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                username,
                display_name,
                role,
                student_number if role == "student" else None,
                salt,
                password_hash,
                now_iso(),
            ),
        )
        created = conn.execute(
            """
            SELECT id, username, display_name, role, student_number, created_at
            FROM users
            WHERE id = ?
            """,
            (cursor.lastrowid,),
        ).fetchone()
        conn.commit()
        conn.close()
        self.send_json({"user": dict(created)}, status=HTTPStatus.CREATED)

    def api_update_user(self, path):
        user = self.require_role("admin")
        if not user:
            return

        try:
            target_user_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "用户编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        data = self.parse_json_body()
        username = (data.get("username") or "").strip()
        display_name = (data.get("display_name") or "").strip()
        role = (data.get("role") or "").strip()
        student_number = (data.get("student_number") or "").strip()
        password = data.get("password") or ""

        if not username or not display_name:
            self.send_json({"error": "用户名和姓名不能为空。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        current = conn.execute(
            "SELECT id, role FROM users WHERE id = ?",
            (target_user_id,),
        ).fetchone()
        if not current:
            conn.close()
            self.send_json({"error": "用户不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if current["role"] == "admin":
            conn.close()
            self.send_json({"error": "管理员账号不支持在此页面修改。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role not in {"teacher", "student"}:
            conn.close()
            self.send_json({"error": "仅支持维护教师和学生账号。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student" and not student_number:
            conn.close()
            self.send_json({"error": "学生账号必须填写学号。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if current["role"] == "teacher" and role != "teacher":
            class_owner = conn.execute(
                "SELECT id FROM classes WHERE teacher_id = ? LIMIT 1",
                (target_user_id,),
            ).fetchone()
            if class_owner:
                conn.close()
                self.send_json({"error": "该教师仍负责班级，请先转移或删除班级。"}, status=HTTPStatus.BAD_REQUEST)
                return

        duplicate_user = conn.execute(
            "SELECT id FROM users WHERE username = ? AND id != ?",
            (username, target_user_id),
        ).fetchone()
        if duplicate_user:
            conn.close()
            self.send_json({"error": "用户名已存在。"}, status=HTTPStatus.BAD_REQUEST)
            return
        if role == "student":
            duplicate_student_number = conn.execute(
                "SELECT id FROM users WHERE student_number = ? AND id != ?",
                (student_number, target_user_id),
            ).fetchone()
            if duplicate_student_number:
                conn.close()
                self.send_json({"error": "学号已存在。"}, status=HTTPStatus.BAD_REQUEST)
                return

        if password:
            salt, password_hash = hash_password(password)
            conn.execute(
                """
                UPDATE users
                SET username = ?, display_name = ?, role = ?, student_number = ?, salt = ?, password_hash = ?
                WHERE id = ?
                """,
                (
                    username,
                    display_name,
                    role,
                    student_number if role == "student" else None,
                    salt,
                    password_hash,
                    target_user_id,
                ),
            )
        else:
            conn.execute(
                """
                UPDATE users
                SET username = ?, display_name = ?, role = ?, student_number = ?
                WHERE id = ?
                """,
                (
                    username,
                    display_name,
                    role,
                    student_number if role == "student" else None,
                    target_user_id,
                ),
            )
        updated = conn.execute(
            """
            SELECT id, username, display_name, role, student_number, created_at
            FROM users
            WHERE id = ?
            """,
            (target_user_id,),
        ).fetchone()
        conn.commit()
        conn.close()
        publish_user_updates(target_user_id)
        self.send_json({"user": dict(updated), "message": "用户信息已更新。"})

    def api_delete_user(self, path):
        user = self.require_role("admin")
        if not user:
            return

        try:
            target_user_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "用户编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        target = conn.execute(
            "SELECT id, role FROM users WHERE id = ?",
            (target_user_id,),
        ).fetchone()
        if not target:
            conn.close()
            self.send_json({"error": "用户不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if target["role"] == "admin":
            conn.close()
            self.send_json({"error": "管理员账号不支持删除。"}, status=HTTPStatus.BAD_REQUEST)
            return

        result = delete_user_record(conn, target_user_id)
        conn.commit()
        conn.close()

        for file_name in result["stored_file_names"]:
            delete_courseware_assets(file_name)

        publish_user_updates(*result["affected_user_ids"])
        self.send_json({"message": "用户已删除。"})

    def api_coursewares(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        class_id = self.get_query_int("class_id")
        if class_id is not None and not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        class_clause, class_params = build_class_scope_clause(user, class_id, "c.class_id")
        rows = conn.execute(
            f"""
            SELECT
                c.id,
                c.title,
                c.course_name,
                c.description,
                c.original_file_name,
                c.stored_file_name,
                c.uploaded_at,
                c.class_id,
                cls.name AS class_name,
                u.display_name AS teacher_name,
                u.id AS teacher_id
            FROM coursewares c
            JOIN classes cls ON cls.id = c.class_id
            JOIN users u ON u.id = c.uploaded_by
            WHERE {class_clause}
            ORDER BY c.uploaded_at DESC
            """,
            class_params,
        ).fetchall()
        conn.close()

        coursewares = []
        for row in rows:
            item = dict(row)
            item["viewer_url"] = build_viewer_url(row["stored_file_name"], row["title"])
            coursewares.append(item)
        self.send_json({"coursewares": coursewares})

    def api_courseware_detail(self, path):
        user = self.require_user()
        if not user:
            return

        try:
            courseware_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "课件编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        row = conn.execute(
            """
            SELECT
                c.id,
                c.title,
                c.course_name,
                c.description,
                c.original_file_name,
                c.stored_file_name,
                c.uploaded_at,
                c.class_id,
                cls.name AS class_name,
                u.display_name AS teacher_name
            FROM coursewares c
            JOIN classes cls ON cls.id = c.class_id
            JOIN users u ON u.id = c.uploaded_by
            WHERE c.id = ?
            """,
            (courseware_id,),
        ).fetchone()

        if not row:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return

        if not self.require_class_access(conn, user, row["class_id"]):
            conn.close()
            return
        conn.close()

        item = dict(row)
        item["viewer_url"] = build_viewer_url(row["stored_file_name"], row["title"])
        self.send_json({"courseware": item})

    def api_create_courseware(self):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        form = self.parse_multipart_form()
        title = (form.get("title", {}).get("value") or "").strip()
        course_name = (form.get("course_name", {}).get("value") or "").strip()
        description = (form.get("description", {}).get("value") or "").strip()
        class_id = form.get("class_id", {}).get("value")
        file_item = form.get("file")

        try:
            class_id = int(class_id)
        except (TypeError, ValueError):
            class_id = None

        if not title or not course_name or class_id is None or file_item is None or not file_item.get("filename"):
            self.send_json({"error": "请填写完整课件信息并选择上传文件。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if not user_can_manage_class(conn, user, class_id):
            conn.close()
            self.send_json({"error": "当前账号无权向该班级上传课件。"}, status=HTTPStatus.FORBIDDEN)
            return
        class_row = conn.execute(
            "SELECT teacher_id FROM classes WHERE id = ?",
            (class_id,),
        ).fetchone()
        if not class_row:
            conn.close()
            self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        uploaded_by = user["id"] if user["role"] == "teacher" else class_row["teacher_id"]

        original_name = os.path.basename(file_item["filename"])
        safe_name = f"{secrets.token_hex(10)}_{original_name}"
        target_path = UPLOAD_DIR / safe_name
        with open(target_path, "wb") as output:
            output.write(file_item["data"])
        if target_path.suffix.lower() in QUICKLOOK_PREVIEW_SUFFIXES:
            if ensure_native_pdf_preview(target_path) is None:
                ensure_quicklook_preview(target_path)
        cursor = conn.execute(
            """
            INSERT INTO coursewares (
                title, course_name, description, original_file_name, stored_file_name, uploaded_by, uploaded_at, class_id
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (title, course_name, description, original_name, safe_name, uploaded_by, now_iso(), class_id),
        )
        courseware_id = cursor.lastrowid
        conn.commit()
        row = conn.execute(
            """
            SELECT
                c.id,
                c.title,
                c.course_name,
                c.description,
                c.original_file_name,
                c.stored_file_name,
                c.uploaded_at,
                c.class_id,
                cls.name AS class_name,
                u.display_name AS teacher_name,
                u.id AS teacher_id
            FROM coursewares c
            JOIN classes cls ON cls.id = c.class_id
            JOIN users u ON u.id = c.uploaded_by
            WHERE c.id = ?
            """,
            (courseware_id,),
        ).fetchone()
        conn.close()

        item = dict(row)
        item["viewer_url"] = build_viewer_url(row["stored_file_name"], row["title"])
        self.send_json({"courseware": item}, status=HTTPStatus.CREATED)

    def api_update_courseware(self, path):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        try:
            courseware_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "课件编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        data = self.parse_json_body()
        title = (data.get("title") or "").strip()
        course_name = (data.get("course_name") or "").strip()
        description = (data.get("description") or "").strip()
        class_id = data.get("class_id")

        if not title or not course_name:
            self.send_json({"error": "标题和课程名称不能为空。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        row = conn.execute(
            "SELECT class_id, uploaded_by FROM coursewares WHERE id = ?",
            (courseware_id,),
        ).fetchone()
        if not row:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not user_can_manage_class(conn, user, row["class_id"]):
            conn.close()
            self.send_json({"error": "当前账号无权修改该课件。"}, status=HTTPStatus.FORBIDDEN)
            return

        new_class_id = row["class_id"]
        uploaded_by = row["uploaded_by"]
        if user["role"] == "admin" and class_id not in (None, ""):
            try:
                new_class_id = int(class_id)
            except (TypeError, ValueError):
                conn.close()
                self.send_json({"error": "班级参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
                return
            target_class = conn.execute(
                "SELECT teacher_id FROM classes WHERE id = ?",
                (new_class_id,),
            ).fetchone()
            if not target_class:
                conn.close()
                self.send_json({"error": "班级不存在。"}, status=HTTPStatus.NOT_FOUND)
                return
            uploaded_by = target_class["teacher_id"]
        elif user["role"] == "teacher":
            uploaded_by = user["id"]

        conn.execute(
            """
            UPDATE coursewares
            SET title = ?, course_name = ?, description = ?, class_id = ?, uploaded_by = ?
            WHERE id = ?
            """,
            (title, course_name, description, new_class_id, uploaded_by, courseware_id),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "课件信息已更新。"})

    def api_delete_courseware(self, path):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return

        try:
            courseware_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "课件编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        row = conn.execute(
            "SELECT stored_file_name, class_id FROM coursewares WHERE id = ?",
            (courseware_id,),
        ).fetchone()
        if not row:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not user_can_manage_class(conn, user, row["class_id"]):
            conn.close()
            self.send_json({"error": "当前账号无权删除该课件。"}, status=HTTPStatus.FORBIDDEN)
            return

        result = delete_courseware_record(conn, courseware_id)
        conn.commit()
        conn.close()

        delete_courseware_assets(row["stored_file_name"])

        publish_user_updates(*result["affected_user_ids"])
        self.send_json({"message": "课件已删除。"})

    def prepare_ai_chat_request(self, user, courseware_id, question):
        if not question:
            self.send_json({"error": "请输入问题内容。"}, status=HTTPStatus.BAD_REQUEST)
            return None

        conn = get_conn()
        courseware = conn.execute(
            """
            SELECT
                c.id,
                c.title,
                c.course_name,
                c.description,
                c.original_file_name,
                c.stored_file_name,
                c.class_id,
                cls.name AS class_name
            FROM coursewares c
            JOIN classes cls ON cls.id = c.class_id
            WHERE c.id = ?
            """,
            (courseware_id,),
        ).fetchone()
        if not courseware:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return None
        if not self.require_class_access(conn, user, courseware["class_id"]):
            conn.close()
            return None

        history_rows = conn.execute(
            """
            SELECT role, content
            FROM ai_chat_messages
            WHERE courseware_id = ? AND user_id = ?
            ORDER BY id DESC
            LIMIT ?
            """,
            (courseware_id, user["id"], MAX_AI_HISTORY_MESSAGES),
        ).fetchall()

        file_path = UPLOAD_DIR / courseware["stored_file_name"]
        courseware_text = extract_courseware_text(file_path)
        context_text = crop_ai_context(courseware_text)
        context_mode = "full" if context_text else "metadata_only"

        history_messages = [
            {"role": row["role"], "content": row["content"]}
            for row in reversed(history_rows)
        ]
        reference_prompt = (
            f"课程名称：{courseware['course_name']}\n"
            f"班级名称：{courseware['class_name']}\n"
            f"课件标题：{courseware['title']}\n"
            f"课件简介：{courseware['description'] or '暂无简介'}\n"
        )
        if context_text:
            reference_prompt += f"\n以下是课件中可解析出的正文内容，请优先基于这些内容回答：\n{context_text}"
        else:
            reference_prompt += (
                "\n当前课件正文暂未成功解析，请仅基于课件标题、课程名、简介和通用教学知识回答，"
                "并在资料不足时明确说明。"
            )

        llm_messages = [
            {
                "role": "system",
                "content": (
                    "你是一名高校课程 AI 助教。"
                    "请使用简洁、准确、友好的中文回答。"
                    "优先依据提供的课件资料作答；如果资料不足，请明确说明，不要编造。"
                    "回答必须尽量结构清晰，优先使用易读的 Markdown 格式。"
                    "请遵守以下格式要求："
                    "先给出一句简短结论；"
                    "再使用 2 到 4 个标题分段；"
                    "段内优先使用项目符号或编号列表；"
                    "避免输出超长连续大段文字；"
                    "标题、段落、列表项之间要有明确换行。"
                ),
            },
            {"role": "system", "content": reference_prompt},
            *history_messages,
            {"role": "user", "content": question},
        ]
        return conn, courseware, llm_messages, context_mode

    def persist_ai_chat_messages(self, conn, courseware_id, user_id, question, answer):
        conn.execute(
            """
            INSERT INTO ai_chat_messages (courseware_id, user_id, role, content, created_at)
            VALUES (?, ?, ?, ?, ?)
            """,
            (courseware_id, user_id, "user", question, now_iso()),
        )
        conn.execute(
            """
            INSERT INTO ai_chat_messages (courseware_id, user_id, role, content, created_at)
            VALUES (?, ?, ?, ?, ?)
            """,
            (courseware_id, user_id, "assistant", answer, now_iso()),
        )
        rows = conn.execute(
            """
            SELECT id, role, content, created_at
            FROM ai_chat_messages
            WHERE courseware_id = ? AND user_id = ?
            ORDER BY id ASC
            """,
            (courseware_id, user_id),
        ).fetchall()
        conn.commit()
        return [dict(row) for row in rows]

    def api_ai_messages(self):
        user = self.require_user()
        if not user:
            return

        courseware_id = self.get_query_int("courseware_id")
        if courseware_id is None:
            self.send_json({"error": "课件参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        courseware = conn.execute(
            "SELECT id, class_id FROM coursewares WHERE id = ?",
            (courseware_id,),
        ).fetchone()
        if not courseware:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not self.require_class_access(conn, user, courseware["class_id"]):
            conn.close()
            return

        rows = conn.execute(
            """
            SELECT id, role, content, created_at
            FROM ai_chat_messages
            WHERE courseware_id = ? AND user_id = ?
            ORDER BY id ASC
            """,
            (courseware_id, user["id"]),
        ).fetchall()
        conn.close()
        self.send_json({"messages": [dict(row) for row in rows]})

    def api_ai_ask(self):
        user = self.require_user()
        if not user:
            return

        data = self.parse_json_body()
        question = (data.get("question") or "").strip()
        courseware_id = data.get("courseware_id")
        try:
            courseware_id = int(courseware_id)
        except (TypeError, ValueError):
            self.send_json({"error": "课件参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        prepared = self.prepare_ai_chat_request(user, courseware_id, question)
        if not prepared:
            return
        conn, _courseware, llm_messages, context_mode = prepared

        try:
            answer = call_bigmodel_chat(llm_messages)
        except RuntimeError as error:
            conn.close()
            status = HTTPStatus.SERVICE_UNAVAILABLE if "BIGMODEL_API_KEY" in str(error) else HTTPStatus.BAD_GATEWAY
            self.send_json({"error": str(error)}, status=status)
            return
        except Exception:
            conn.close()
            self.send_json({"error": "AI 服务暂时不可用，请稍后重试。"}, status=HTTPStatus.BAD_GATEWAY)
            return

        rows = self.persist_ai_chat_messages(conn, courseware_id, user["id"], question, answer)
        conn.close()
        self.send_json(
            {
                "answer": answer,
                "messages": rows,
                "context_mode": context_mode,
            },
            status=HTTPStatus.CREATED,
        )

    def api_clear_ai_messages(self):
        user = self.require_user()
        if not user:
            return

        courseware_id = self.get_query_int("courseware_id")
        if courseware_id is None:
            self.send_json({"error": "课件参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        courseware = conn.execute(
            "SELECT id, class_id FROM coursewares WHERE id = ?",
            (courseware_id,),
        ).fetchone()
        if not courseware:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not self.require_class_access(conn, user, courseware["class_id"]):
            conn.close()
            return

        conn.execute(
            "DELETE FROM ai_chat_messages WHERE courseware_id = ? AND user_id = ?",
            (courseware_id, user["id"]),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "问答记录已清空。"})

    def api_evaluations(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        class_id = self.get_query_int("class_id")
        if class_id is not None and not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        class_clause, class_params = build_class_scope_clause(user, class_id, "c.class_id")
        if user["role"] == "teacher":
            rows = conn.execute(
                f"""
                SELECT
                    e.id,
                    e.helpfulness AS difficulty,
                    e.usability AS readability,
                    COALESCE(e.suitability, e.usability) AS suitability,
                    COALESCE(e.practicality, e.helpfulness) AS practicality,
                    e.suggestion,
                    e.created_at,
                    c.class_id,
                    cls.name AS class_name,
                    c.title AS courseware_title,
                    c.course_name,
                    u.display_name AS student_name
                FROM evaluations e
                JOIN coursewares c ON c.id = e.courseware_id
                JOIN classes cls ON cls.id = c.class_id
                JOIN users u ON u.id = e.student_id
                WHERE {class_clause}
                ORDER BY e.created_at DESC
                """,
                class_params,
            ).fetchall()
        else:
            rows = conn.execute(
                f"""
                SELECT
                    e.id,
                    e.helpfulness AS difficulty,
                    e.usability AS readability,
                    COALESCE(e.suitability, e.usability) AS suitability,
                    COALESCE(e.practicality, e.helpfulness) AS practicality,
                    e.suggestion,
                    e.created_at,
                    c.class_id,
                    cls.name AS class_name,
                    c.title AS courseware_title,
                    c.course_name
                FROM evaluations e
                JOIN coursewares c ON c.id = e.courseware_id
                JOIN classes cls ON cls.id = c.class_id
                WHERE e.student_id = ? AND {class_clause}
                ORDER BY e.created_at DESC
                """,
                [user["id"], *class_params],
            ).fetchall()
        conn.close()
        self.send_json({"evaluations": [dict(row) for row in rows]})

    def api_create_evaluation(self):
        user = self.require_role("student")
        if not user:
            return

        data = self.parse_json_body()
        courseware_id = data.get("courseware_id")
        difficulty = data.get("difficulty")
        readability = data.get("readability")
        suitability = data.get("suitability")
        practicality = data.get("practicality")
        suggestion = (data.get("suggestion") or "").strip()

        try:
            courseware_id = int(courseware_id)
            difficulty = int(difficulty)
            readability = int(readability)
            suitability = int(suitability)
            practicality = int(practicality)
        except (TypeError, ValueError):
            self.send_json({"error": "问卷参数格式不正确。"}, status=HTTPStatus.BAD_REQUEST)
            return

        if not all(1 <= score <= 5 for score in [difficulty, readability, suitability, practicality]):
            self.send_json({"error": "评分范围应在 1 到 5 之间。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        courseware = conn.execute("SELECT id, class_id FROM coursewares WHERE id = ?", (courseware_id,)).fetchone()
        if not courseware:
            conn.close()
            self.send_json({"error": "课件不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not self.require_class_access(conn, user, courseware["class_id"]):
            conn.close()
            return

        exists = conn.execute(
            "SELECT id FROM evaluations WHERE courseware_id = ? AND student_id = ?",
            (courseware_id, user["id"]),
        ).fetchone()
        if exists:
            conn.close()
            self.send_json({"error": "同一课件只允许提交一次问卷。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn.execute(
            """
            INSERT INTO evaluations (
                courseware_id, student_id, helpfulness, usability, suitability, practicality, suggestion, created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                courseware_id,
                user["id"],
                difficulty,
                readability,
                suitability,
                practicality,
                suggestion,
                now_iso(),
            ),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "问卷提交成功。"}, status=HTTPStatus.CREATED)

    def api_discussions(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        class_id = self.get_query_int("class_id")
        if class_id is not None and not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        class_clause, class_params = build_class_scope_clause(user, class_id, "d.class_id")
        discussion_rows = conn.execute(
            f"""
            SELECT
                d.id,
                d.title,
                d.body,
                d.created_at,
                d.class_id,
                cls.name AS class_name,
                u.display_name AS author_name,
                u.role AS author_role
            FROM discussions d
            JOIN classes cls ON cls.id = d.class_id
            JOIN users u ON u.id = d.author_id
            WHERE {class_clause}
            ORDER BY d.created_at DESC
            """,
            class_params,
        ).fetchall()
        reply_rows = conn.execute(
            """
            SELECT
                r.id,
                r.discussion_id,
                r.body,
                r.created_at,
                u.display_name AS author_name,
                u.role AS author_role
            FROM discussion_replies r
            JOIN users u ON u.id = r.author_id
            ORDER BY r.created_at ASC
            """
        ).fetchall()
        conn.close()

        replies_map = {}
        for row in reply_rows:
            replies_map.setdefault(row["discussion_id"], []).append(dict(row))

        discussions = []
        for row in discussion_rows:
            item = dict(row)
            item["replies"] = replies_map.get(row["id"], [])
            discussions.append(item)

        self.send_json({"discussions": discussions})

    def api_create_discussion(self):
        user = self.require_user()
        if not user:
            return

        data = self.parse_json_body()
        title = (data.get("title") or "").strip()
        body = (data.get("body") or "").strip()
        class_id = data.get("class_id")
        try:
            class_id = int(class_id)
        except (TypeError, ValueError):
            class_id = None

        if not title or not body or class_id is None:
            self.send_json({"error": "请输入帖子标题和内容。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        conn.execute(
            "INSERT INTO discussions (title, body, author_id, created_at, class_id) VALUES (?, ?, ?, ?, ?)",
            (title, body, user["id"], now_iso(), class_id),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "帖子发布成功。"}, status=HTTPStatus.CREATED)

    def api_create_discussion_reply(self, path):
        user = self.require_user()
        if not user:
            return

        try:
            discussion_id = int(path.split("/")[-2])
        except ValueError:
            self.send_json({"error": "帖子编号不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        data = self.parse_json_body()
        body = (data.get("body") or "").strip()
        if not body:
            self.send_json({"error": "回复内容不能为空。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        discussion = conn.execute("SELECT id, class_id FROM discussions WHERE id = ?", (discussion_id,)).fetchone()
        if not discussion:
            conn.close()
            self.send_json({"error": "帖子不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not self.require_class_access(conn, user, discussion["class_id"]):
            conn.close()
            return

        conn.execute(
            """
            INSERT INTO discussion_replies (discussion_id, body, author_id, created_at)
            VALUES (?, ?, ?, ?)
            """,
            (discussion_id, body, user["id"], now_iso()),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "回复已发布。"}, status=HTTPStatus.CREATED)

    def api_conversations(self):
        user = self.require_user()
        if not user:
            return

        conn = get_conn()
        rows = conn.execute(
            """
            SELECT
                ct.id AS thread_id,
                ct.last_message_at,
                other_user.id AS user_id,
                other_user.display_name,
                other_user.role,
                latest_message.body AS last_message,
                (
                    SELECT COUNT(*)
                    FROM messages unread_message
                    WHERE unread_message.thread_id = ct.id
                      AND unread_message.receiver_id = ?
                      AND unread_message.is_read = 0
                ) AS unread_count
            FROM conversation_members current_member
            JOIN conversation_threads ct ON ct.id = current_member.thread_id
            JOIN conversation_members other_member
              ON other_member.thread_id = ct.id AND other_member.user_id != current_member.user_id
            JOIN users other_user ON other_user.id = other_member.user_id
            LEFT JOIN messages latest_message
              ON latest_message.id = (
                  SELECT id
                  FROM messages
                  WHERE thread_id = ct.id
                  ORDER BY created_at DESC, id DESC
                  LIMIT 1
              )
            WHERE current_member.user_id = ?
              AND current_member.visible = 1
            ORDER BY ct.last_message_at DESC, ct.id DESC
            """,
            (user["id"], user["id"]),
        ).fetchall()
        conn.close()

        conversations = [
            {
                "thread_id": row["thread_id"],
                "user": {
                    "id": row["user_id"],
                    "display_name": row["display_name"],
                    "role": row["role"],
                },
                "last_message": row["last_message"] or "",
                "last_message_at": row["last_message_at"],
                "unread_count": row["unread_count"],
            }
            for row in rows
        ]
        self.send_json({"conversations": conversations})

    def api_message_contacts(self):
        user = self.require_user()
        if not user:
            return

        class_id = self.get_query_int("class_id")
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return

        rows = conn.execute(
            """
            SELECT DISTINCT u.id, u.display_name, u.username, u.role
            FROM class_members cm
            JOIN users u ON u.id = cm.user_id
            WHERE cm.class_id = ? AND u.id != ?
            ORDER BY CASE u.role WHEN 'teacher' THEN 0 ELSE 1 END, u.display_name ASC
            """,
            (class_id, user["id"]),
        ).fetchall()

        contacts = []
        for row in rows:
            thread_id = find_thread_id(conn, user["id"], row["id"])
            current_member = None
            if thread_id:
                current_member = conn.execute(
                    """
                    SELECT visible
                    FROM conversation_members
                    WHERE thread_id = ? AND user_id = ?
                    """,
                    (thread_id, user["id"]),
                ).fetchone()
            contacts.append(
                {
                    "id": row["id"],
                    "display_name": row["display_name"],
                    "username": row["username"],
                    "role": row["role"],
                    "has_conversation": bool(current_member and current_member["visible"]),
                }
            )
        conn.close()
        self.send_json({"contacts": contacts})

    def api_message_events(self):
        user = self.require_user()
        if not user:
            return

        cursor = self.get_query_int("cursor")
        if cursor is None:
            cursor = 0

        latest_cursor = get_user_sync_cursor(user["id"])
        changed = latest_cursor > cursor
        if not changed:
            latest_cursor = wait_for_user_update(user["id"], cursor)
            changed = latest_cursor > cursor

        self.send_json({"changed": changed, "cursor": latest_cursor})

    def api_create_conversation(self):
        user = self.require_user()
        if not user:
            return

        data = self.parse_json_body()
        contact_id = data.get("contact_id")
        try:
            contact_id = int(contact_id)
        except (TypeError, ValueError):
            self.send_json({"error": "联系人参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        contact = conn.execute(
            "SELECT id FROM users WHERE id = ?",
            (contact_id,),
        ).fetchone()
        if not contact:
            conn.close()
            self.send_json({"error": "联系人不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not users_share_class(conn, user["id"], contact_id):
            conn.close()
            self.send_json({"error": "仅可与同班联系人创建会话。"}, status=HTTPStatus.FORBIDDEN)
            return

        thread_id = get_or_create_thread(conn, user["id"], contact_id, visible_for_a=1, visible_for_b=0)
        conn.execute(
            """
            UPDATE conversation_members
            SET visible = 1
            WHERE thread_id = ? AND user_id = ?
            """,
            (thread_id, user["id"]),
        )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"])
        self.send_json({"message": "会话已加入列表。"}, status=HTTPStatus.CREATED)

    def api_delete_conversation(self, path):
        user = self.require_user()
        if not user:
            return

        try:
            contact_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "联系人参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        thread_id = find_thread_id(conn, user["id"], contact_id)
        if thread_id:
            conn.execute(
                """
                UPDATE conversation_members
                SET visible = 0
                WHERE thread_id = ? AND user_id = ?
                """,
                (thread_id, user["id"]),
            )
            conn.commit()
        conn.close()
        publish_user_updates(user["id"])
        self.send_json({"message": "会话已从列表移除。"})

    def api_message_thread(self, path):
        user = self.require_user()
        if not user:
            return

        try:
            other_id = int(path.split("/")[-1])
        except ValueError:
            self.send_json({"error": "会话对象不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        other = conn.execute(
            "SELECT id, username, display_name, role, created_at FROM users WHERE id = ?",
            (other_id,),
        ).fetchone()
        if not other:
            conn.close()
            self.send_json({"error": "用户不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not users_share_class(conn, user["id"], other_id):
            conn.close()
            self.send_json({"error": "仅可查看同班联系人会话。"}, status=HTTPStatus.FORBIDDEN)
            return

        thread_id = get_or_create_thread(conn, user["id"], other_id, visible_for_a=1, visible_for_b=0)

        cursor = conn.execute(
            """
            UPDATE messages
            SET is_read = 1
            WHERE thread_id = ? AND sender_id = ? AND receiver_id = ? AND is_read = 0
            """,
            (thread_id, other_id, user["id"]),
        )
        rows = conn.execute(
            """
            SELECT id, thread_id, sender_id, receiver_id, body, is_read, created_at
            FROM messages
            WHERE thread_id = ?
            ORDER BY created_at ASC, id ASC
            """,
            (thread_id,),
        ).fetchall()
        conn.commit()
        conn.close()
        if cursor.rowcount:
            publish_user_updates(user["id"], other_id)
        self.send_json({"thread_id": thread_id, "other_user": dict(other), "messages": [dict(row) for row in rows]})

    def api_send_message(self):
        user = self.require_user()
        if not user:
            return

        data = self.parse_json_body()
        receiver_id = data.get("receiver_id")
        body = (data.get("body") or "").strip()
        if not body:
            self.send_json({"error": "消息内容不能为空。"}, status=HTTPStatus.BAD_REQUEST)
            return

        try:
            receiver_id = int(receiver_id)
        except (TypeError, ValueError):
            self.send_json({"error": "接收人参数不合法。"}, status=HTTPStatus.BAD_REQUEST)
            return

        conn = get_conn()
        receiver = conn.execute("SELECT id FROM users WHERE id = ?", (receiver_id,)).fetchone()
        if not receiver:
            conn.close()
            self.send_json({"error": "接收人不存在。"}, status=HTTPStatus.NOT_FOUND)
            return
        if not users_share_class(conn, user["id"], receiver_id):
            conn.close()
            self.send_json({"error": "仅可向同班联系人发送私信。"}, status=HTTPStatus.FORBIDDEN)
            return

        thread_id = get_or_create_thread(conn, user["id"], receiver_id, visible_for_a=1, visible_for_b=1)
        conn.execute(
            """
            UPDATE conversation_members
            SET visible = 1
            WHERE thread_id = ? AND user_id IN (?, ?)
            """,
            (thread_id, user["id"], receiver_id),
        )

        conn.execute(
            """
            INSERT INTO messages (sender_id, receiver_id, body, is_read, created_at, thread_id)
            VALUES (?, ?, ?, 0, ?, ?)
            """,
            (user["id"], receiver_id, body, now_iso(), thread_id),
        )
        conn.execute(
            "UPDATE conversation_threads SET last_message_at = ? WHERE id = ?",
            (now_iso(), thread_id),
        )
        conn.commit()
        conn.close()
        publish_user_updates(user["id"], receiver_id)
        self.send_json({"message": "消息发送成功。"}, status=HTTPStatus.CREATED)

    # ------------------------------------------------------------------
    # RAG endpoints
    # ------------------------------------------------------------------

    def api_rag_status(self):
        user = self.require_user()
        if not user:
            return
        class_id = self.get_query_int("class_id")
        if class_id is None:
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        conn.close()
        if not RAG_AVAILABLE:
            self.send_json({"available": False, "building": False, "indexed": False, "chunk_count": 0})
            return
        status = _rag_module.get_index_status(class_id)
        status["available"] = True
        self.send_json(status)

    def api_rag_index(self):
        user = self.require_role(("teacher", "admin"))
        if not user:
            return
        if not RAG_AVAILABLE:
            self.send_json({"error": "RAG 模块不可用，请检查依赖安装。"}, status=HTTPStatus.SERVICE_UNAVAILABLE)
            return
        data = self.parse_json_body()
        try:
            class_id = int(data.get("class_id"))
        except (TypeError, ValueError):
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        rows = conn.execute(
            """
            SELECT id, title, stored_file_name, course_name
            FROM coursewares
            WHERE class_id = ?
            """,
            (class_id,),
        ).fetchall()
        conn.close()
        courseware_list = [dict(row) for row in rows]
        _rag_module.build_class_index_async(class_id, courseware_list, UPLOAD_DIR, extract_courseware_text)
        self.send_json({"message": "索引构建任务已启动，请稍后刷新状态。"})

    def api_rag_ask(self):
        user = self.require_user()
        if not user:
            return
        if not RAG_AVAILABLE:
            self.send_json({"error": "RAG 模块不可用，请检查依赖安装。"}, status=HTTPStatus.SERVICE_UNAVAILABLE)
            return
        data = self.parse_json_body()
        question = (data.get("question") or "").strip()
        if not question:
            self.send_json({"error": "请输入问题。"}, status=HTTPStatus.BAD_REQUEST)
            return
        try:
            class_id = int(data.get("class_id"))
        except (TypeError, ValueError):
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        # Load recent history
        history_rows = conn.execute(
            """
            SELECT role, content
            FROM rag_chat_messages
            WHERE class_id = ? AND user_id = ?
            ORDER BY id DESC
            LIMIT ?
            """,
            (class_id, user["id"], MAX_AI_HISTORY_MESSAGES),
        ).fetchall()
        history = [{"role": row["role"], "content": row["content"]} for row in reversed(history_rows)]
        try:
            result = _rag_module.rag_ask(class_id, question, history_messages=history)
        except RuntimeError as exc:
            conn.close()
            self.send_json({"error": str(exc)}, status=HTTPStatus.BAD_GATEWAY)
            return
        except Exception:
            conn.close()
            self.send_json({"error": "AI 服务暂时不可用，请稍后重试。"}, status=HTTPStatus.BAD_GATEWAY)
            return
        answer = result["answer"]
        sources = result["sources"]
        sources_json = json.dumps(sources, ensure_ascii=False)
        now = now_iso()
        conn.execute(
            "INSERT INTO rag_chat_messages (class_id, user_id, role, content, sources, created_at) VALUES (?, ?, ?, ?, ?, ?)",
            (class_id, user["id"], "user", question, "[]", now),
        )
        conn.execute(
            "INSERT INTO rag_chat_messages (class_id, user_id, role, content, sources, created_at) VALUES (?, ?, ?, ?, ?, ?)",
            (class_id, user["id"], "assistant", answer, sources_json, now),
        )
        rows = conn.execute(
            """
            SELECT id, role, content, sources, created_at
            FROM rag_chat_messages
            WHERE class_id = ? AND user_id = ?
            ORDER BY id ASC
            """,
            (class_id, user["id"]),
        ).fetchall()
        conn.commit()
        conn.close()
        messages = []
        for row in rows:
            item = dict(row)
            try:
                item["sources"] = json.loads(item["sources"])
            except Exception:
                item["sources"] = []
            messages.append(item)
        self.send_json({"answer": answer, "sources": sources, "messages": messages}, status=HTTPStatus.CREATED)

    def api_rag_messages(self):
        user = self.require_user()
        if not user:
            return
        class_id = self.get_query_int("class_id")
        if class_id is None:
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        rows = conn.execute(
            """
            SELECT id, role, content, sources, created_at
            FROM rag_chat_messages
            WHERE class_id = ? AND user_id = ?
            ORDER BY id ASC
            """,
            (class_id, user["id"]),
        ).fetchall()
        conn.close()
        messages = []
        for row in rows:
            item = dict(row)
            try:
                item["sources"] = json.loads(item["sources"])
            except Exception:
                item["sources"] = []
            messages.append(item)
        self.send_json({"messages": messages})

    def api_rag_clear_messages(self):
        user = self.require_user()
        if not user:
            return
        class_id = self.get_query_int("class_id")
        if class_id is None:
            self.send_json({"error": "请先选择班级。"}, status=HTTPStatus.BAD_REQUEST)
            return
        conn = get_conn()
        if not self.require_class_access(conn, user, class_id):
            conn.close()
            return
        conn.execute(
            "DELETE FROM rag_chat_messages WHERE class_id = ? AND user_id = ?",
            (class_id, user["id"]),
        )
        conn.commit()
        conn.close()
        self.send_json({"message": "问答记录已清空。"})


class AppServer(ThreadingHTTPServer):
    allow_reuse_address = True
    daemon_threads = True


def run():
    parser = argparse.ArgumentParser(description="Run the AI Tutor System server.")
    parser.add_argument("--host", default="127.0.0.1", help="Host to bind, default is 127.0.0.1")
    parser.add_argument("--port", type=int, default=8000, help="Port to bind, default is 8000")
    args = parser.parse_args()

    init_db()
    try:
        server = AppServer((args.host, args.port), AppHandler)
    except OSError as error:
        if error.errno == 48:
            print(
                f"Port {args.port} is already in use. "
                f"Stop the existing process or run: python3 server.py --port {args.port + 1}"
            )
            return
        raise

    print(f"AI Tutor System running at http://{args.host}:{args.port}")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nServer stopped.")


if __name__ == "__main__":
    run()
